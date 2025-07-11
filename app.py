import os
import json
import logging
from datetime import datetime
from pathlib import Path
from flask import Flask, render_template, request, jsonify, send_file, send_from_directory
from flask_cors import CORS
import chromadb
from chromadb.utils import embedding_functions
from openai import OpenAI
import tempfile
import shutil
from fpdf import FPDF
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import glob
from dotenv import load_dotenv
from rag_system import EnhancedDentalRAG
from database_manager import PracticeDatabase
from datetime import datetime, timedelta
import sqlite3
import html2text

# Load environment variables from .env file
load_dotenv()

app = Flask(__name__)
CORS(app)

# Initialize OpenAI client
api_key = os.getenv('OPENAI_API_KEY')
if not api_key:
    raise ValueError("OPENAI_API_KEY not found in environment variables. Please set it in your .env file.")

client = OpenAI(api_key=api_key)

# Initialize enhanced RAG system
print("🚀 Initializing Enhanced Dental RAG System...")
rag_system = EnhancedDentalRAG()

# Initialize practice database
print("🏥 Initializing Practice Management Database...")
practice_db = PracticeDatabase()

# Auto-index on startup
print("📚 Auto-indexing dental knowledge...")
rag_system.reindex_all()

class SpecializedLLM:
    """Specialized LLM instance for each tab with focused context and prompts"""
    
    def __init__(self, tab_name: str, system_prompt: str, rag_system: EnhancedDentalRAG):
        self.tab_name = tab_name
        self.base_system_prompt = system_prompt
        self.rag_system = rag_system
        self.chat_history = []
        
    def get_specialized_context(self, user_message: str) -> tuple:
        """Get context specifically relevant to this tab"""
        if self.tab_name == 'dental-brain':
            # For dental-brain: cases + general knowledge
            rag_results = self.rag_system.search_combined(
                user_message, 
                case_results=3,  # Increased from 1 to 3
                knowledge_results=2  # Reduced from 3
            )
        elif self.tab_name == 'swiss-law':
            # For Swiss law: only legal knowledge
            rag_results = self.rag_system.search_knowledge(user_message, n_results=2)  # Reduced from 4
            rag_results = {'cases': [], 'knowledge': rag_results, 'total_results': len(rag_results)}
        elif self.tab_name == 'invisalign':
            # For Invisalign: orthodontic knowledge
            rag_results = self.rag_system.search_knowledge(user_message, n_results=2)  # Reduced from 4
            rag_results = {'cases': [], 'knowledge': rag_results, 'total_results': len(rag_results)}
        elif self.tab_name == 'office-knowledge':
            # For office procedures: internal knowledge
            rag_results = self.rag_system.search_knowledge(user_message, n_results=2)  # Reduced from 4
            rag_results = {'cases': [], 'knowledge': rag_results, 'total_results': len(rag_results)}
        elif self.tab_name == 'insurance':
            # For insurance: billing and TARMED knowledge
            rag_results = self.rag_system.search_knowledge(user_message, n_results=2)  # Reduced from 4
            rag_results = {'cases': [], 'knowledge': rag_results, 'total_results': len(rag_results)}
        elif self.tab_name == 'patient-comm':
            # For patient communication: communication knowledge
            rag_results = self.rag_system.search_knowledge(user_message, n_results=2)  # Reduced from 4
            rag_results = {'cases': [], 'knowledge': rag_results, 'total_results': len(rag_results)}
        elif self.tab_name == 'emergency':
            # For emergency: emergency protocols
            rag_results = self.rag_system.search_knowledge(user_message, n_results=2)  # Reduced from 4
            rag_results = {'cases': [], 'knowledge': rag_results, 'total_results': len(rag_results)}
        elif self.tab_name == 'patient-education':
            # For patient education: communication and educational knowledge
            rag_results = self.rag_system.search_knowledge(user_message, n_results=2)
            rag_results = {'cases': [], 'knowledge': rag_results, 'total_results': len(rag_results)}
        else:
            rag_results = {'cases': [], 'knowledge': [], 'total_results': 0}
        
        return self.build_focused_context(rag_results)
    
    def build_focused_context(self, rag_results: dict) -> tuple:
        """Build context focused on this tab's domain"""
        context_parts = []
        references = []
        
        # Add cases if relevant (mainly for dental-brain)
        if rag_results['cases']:
            context_parts.append("=== CAS SIMILAIRES ===")
            for i, case in enumerate(rag_results['cases'], 1):
                # Skip the problematic metadata entry
                if case.get('id') == 'metadata':
                    continue
                    
                consultation = case['metadata'].get('consultation', 'Non spécifiée')[:80]  # Reduced from 100
                context_parts.append(f"Cas {i}: {consultation}")
                
                content_preview = case['content'][:150] + "..." if len(case['content']) > 150 else case['content']  # Reduced from 200
                context_parts.append(f"Détails: {content_preview}")
                
                references.append({
                    'id': case.get('id', f"case_{i}"),
                    'title': f"Cas clinique {i}",
                    'description': consultation[:40] + "...",  # Reduced from 50
                    'similarity': case['similarity'],
                    'type': 'case',
                    'content': case['content']
                })
        
        # Add specialized knowledge
        if rag_results['knowledge']:
            context_parts.append(f"=== CONNAISSANCES ({self.tab_name.upper()}) ===")
            for i, knowledge in enumerate(rag_results['knowledge'], 1):
                # Filter knowledge by relevance to this tab
                if self.is_knowledge_relevant(knowledge):
                    title = knowledge['metadata'].get('title', 'N/A')[:50]  # Limit title length
                    context_parts.append(f"Connaissance {i}: {title}")
                    
                    # Truncate content appropriately
                    content = knowledge['content'][:200] + "..." if len(knowledge['content']) > 200 else knowledge['content']  # Reduced from 300
                    context_parts.append(f"Contenu: {content}")
                    
                    references.append({
                        'id': knowledge.get('id', f"knowledge_{i}"),
                        'title': title,
                        'description': knowledge['content'][:50] + "..." if len(knowledge['content']) > 50 else knowledge['content'],  # Reduced from 75
                        'similarity': knowledge['similarity'],
                        'type': 'knowledge',
                        'category': knowledge['metadata'].get('category', 'Non spécifiée'),
                        'content': knowledge['content']
                    })
        
        # Build final context
        if context_parts:
            context = "\n".join(context_parts)
        else:
            context = f"Aucun contexte spécifique trouvé pour {self.tab_name}."
        
        return context, references
    
    def is_knowledge_relevant(self, knowledge: dict) -> bool:
        """Check if knowledge is relevant to this tab"""
        metadata = knowledge.get('metadata', {})
        category = metadata.get('category', '').lower()
        knowledge_type = metadata.get('type', '').lower()
        
        # Tab-specific relevance filters
        if self.tab_name == 'swiss-law':
            return 'law' in category or 'legal' in category or 'suisse' in category
        elif self.tab_name == 'invisalign':
            return 'invisalign' in category or 'orthodon' in category
        elif self.tab_name == 'office-knowledge':
            return 'office' in category or 'procedure' in category or 'cabinet' in category
        elif self.tab_name == 'insurance':
            return 'insurance' in category or 'tarmed' in category or 'billing' in category
        elif self.tab_name == 'patient-comm':
            return 'communication' in category or 'patient' in category
        elif self.tab_name == 'emergency':
            return 'emergency' in category or 'urgence' in category
        elif self.tab_name == 'patient-education':
            return 'communication' in category or 'patient' in category or 'education' in category
        else:
            return True  # For dental-brain, include all knowledge
    
    def generate_response(self, user_message: str) -> dict:
        """Generate response using this specialized LLM"""
        try:
            # Get specialized context
            context, references = self.get_specialized_context(user_message)
            
            # Build complete system prompt (removed terminology injection)
            system_prompt = self.base_system_prompt
            if context:
                system_prompt += f"\n\n=== CONTEXTE SPÉCIALISÉ ===\n{context}"
            
            # Prepare messages
            messages = [
                {"role": "system", "content": system_prompt},
                *self.chat_history,
                {"role": "user", "content": user_message}
            ]
            
            # Call OpenAI API
            response = client.chat.completions.create(
                model="gpt-4",
                messages=messages,
                max_tokens=1500,  # Reduced from 2000
                temperature=0.7
            )
            
            ai_response = response.choices[0].message.content
            
            # Update chat history
            self.chat_history.append({"role": "user", "content": user_message})
            self.chat_history.append({"role": "assistant", "content": ai_response})
            
            # Keep chat history manageable (last 6 messages instead of 10)
            if len(self.chat_history) > 6:
                self.chat_history = self.chat_history[-6:]
            
            return {
                'success': True,
                'response': ai_response,
                'references': references,
                'context_info': {
                    'context_length': len(context),
                    'references_count': len(references),
                    'tab': self.tab_name
                }
            }
            
        except Exception as e:
            return {
                'success': False,
                'error': str(e)
            }

# Initialize specialized LLM instances
specialized_llms = {}

def get_specialized_system_prompts():
    """Get system prompts for each specialized LLM"""
    return {
        'dental-brain': """Tu es un assistant dentaire spécialisé dans la planification de traitements.

EXPERTISE:
- Planification de séquences de traitement dentaire
- Analyse de consultations dentaires
- Génération de plans de traitement au format JSON

FORMAT DE RÉPONSE pour les plans de traitement:
{
  "consultation_text": "Répète la demande",
  "patient_info": {
    "dent": "Numéro de dent",
    "diagnostic": "Diagnostic principal"
  },
  "treatment_sequence": [
    {
      "rdv": 1,
      "traitement": "Description",
      "duree": "Durée estimée",
      "delai": "Délai avant prochain RDV",
      "dr": "Dr. [Nom]",
      "remarque": "Notes"
    }
  ]
}

Retournez le JSON directement, sans balises markdown.""",

        'swiss-law': """Tu es un expert en droit dentaire suisse. Tu maîtrises parfaitement :
- La loi fédérale sur les professions médicales (LPMéd)
- Les réglementations cantonales spécifiques
- Les obligations de formation continue
- Les responsabilités professionnelles
- La protection des données patients
- Les aspects de facturation et remboursement

Réponds de manière précise et cite les articles de loi pertinents quand possible.""",

        'invisalign': """Tu es un orthodontiste expert en Invisalign. Tu maîtrises :
- La sélection des cas appropriés pour Invisalign
- La planification ClinCheck avancée
- Les techniques d'attachements et IPR
- La gestion des mouvements complexes
- Le suivi et les refinements
- Les protocoles de rétention

Donne des conseils pratiques et techniques basés sur l'expérience clinique.""",

        'office-knowledge': """Tu es l'assistant du cabinet dentaire. Tu connais :
- Les procédures administratives internes
- Les protocoles de prise de rendez-vous
- Les questions fréquentes du personnel
- Les procédures d'urgence
- L'organisation du cabinet
- Les relations avec les patients

Réponds de manière pratique et orientée solutions.""",

        'insurance': """Tu es un expert en assurances dentaires suisses. Tu maîtrises :
- Les codes TARMED et leur application
- Les remboursements LAMal et LCA
- Les assurances complémentaires
- La facturation privée
- Les procédures de remboursement
- Les litiges avec les assurances

Donne des conseils précis sur la facturation et les remboursements.""",

        'patient-comm': """Tu es un expert en communication patient. Tu excelles dans :
- L'explication des traitements dentaires
- La gestion des anxiétés et peurs
- La motivation à l'hygiène bucco-dentaire
- L'éducation thérapeutique
- La communication empathique
- Les techniques de persuasion bienveillante

Aide à créer des messages clairs et rassurants pour les patients.""",

        'emergency': """Tu es un expert en urgences dentaires. Tu maîtrises :
- Le diagnostic différentiel des douleurs
- Les protocoles d'urgence
- La gestion de la douleur
- Les prescriptions d'urgence
- Les traumatismes dentaires
- L'orientation vers les spécialistes

Donne des conseils rapides et sûrs pour les situations d'urgence.""",

        'patient-education': """Tu es un assistant spécialisé dans l'éducation des patients dentaires.

EXPERTISE:
- Vulgarisation des traitements dentaires
- Explication des procédures en langage simple
- Motivation et rassurance des patients
- Instructions de soins et prévention

MISSION:
Créer des documents éducatifs pour les patients qui expliquent:
- Pourquoi le traitement est nécessaire
- Comment se déroule chaque étape
- Les bénéfices attendus
- Les soins post-traitement
- Les conseils de prévention

STYLE:
- Langage simple et accessible
- Ton rassurant et professionnel
- Structure claire avec titres et sections
- Éviter le jargon médical
- Utiliser des analogies compréhensibles

FORMAT DE RÉPONSE:
Retournez un document structuré avec:
- Titre accrocheur
- Introduction rassurante
- Explication du diagnostic
- Description du traitement étape par étape
- Bénéfices et résultats attendus
- Instructions post-traitement
- Conseils de prévention
- Contact pour questions""",

        'schedule': """Tu es un assistant intelligent de planification dentaire. Tu maîtrises :
- L'optimisation des plannings de cabinet dentaire
- La gestion des urgences et reprogrammations
- L'analyse des préférences praticiens et patients
- Les règles de bonnes pratiques de programmation
- La gestion des conflits d'horaires
- L'optimisation des créneaux libres

Ton rôle est d'aider le dentiste à gérer son planning de manière optimale. Tu peux :
- Reprogrammer des rendez-vous selon les contraintes
- Proposer des créneaux alternatifs
- Optimiser l'organisation des journées
- Gérer les urgences en trouvant des solutions
- Analyser le planning pour suggérer des améliorations

Tu dois toujours prendre en compte :
- Les préférences du praticien (horaires, types de traitements)
- Les contraintes patients (disponibilités, urgences)
- Les durées des traitements
- Les règles de bonnes pratiques (espacement, récupération)
- L'efficacité globale du cabinet

Réponds de manière pratique avec des solutions concrètes et actionables."""
    }

def initialize_specialized_llms():
    """Initialize all specialized LLM instances"""
    global specialized_llms
    
    system_prompts = get_specialized_system_prompts()
    
    for tab_name, system_prompt in system_prompts.items():
        specialized_llms[tab_name] = SpecializedLLM(tab_name, system_prompt, rag_system)
    
    print("🤖 Initialized specialized LLMs:")
    for tab_name in specialized_llms.keys():
        print(f"   - {tab_name}")

# Initialize specialized LLMs on startup
initialize_specialized_llms()

def parse_duration(duration_str):
    """Parse duration string to minutes"""
    if 'h' in duration_str:
        hours = float(duration_str.replace('h', '').strip())
        return int(hours * 60)
    elif 'min' in duration_str:
        return int(duration_str.replace('min', '').strip())
    else:
        return 60  # Default to 1 hour

# === PRACTICE MANAGEMENT ENDPOINTS ===

@app.route('/api/patients', methods=['GET', 'POST'])
def manage_patients():
    """Manage patients - GET to retrieve, POST to create"""
    if request.method == 'GET':
        search_term = request.args.get('search', '')
        try:
            patients = practice_db.get_patients(search_term if search_term else None)
            return jsonify({
                'success': True,
                'patients': patients,
                'count': len(patients)
            })
        except Exception as e:
            return jsonify({
                'success': False,
                'error': str(e)
            }), 500
    
    elif request.method == 'POST':
        try:
            patient_data = request.get_json()
            
            # Validate required fields
            if not patient_data.get('first_name') or not patient_data.get('last_name'):
                return jsonify({
                    'success': False,
                    'error': 'First name and last name are required'
                }), 400
            
            patient_id = practice_db.create_patient(**patient_data)
            return jsonify({
                'success': True,
                'patient_id': patient_id,
                'message': 'Patient created successfully'
            })
        except Exception as e:
            return jsonify({
                'success': False,
                'error': str(e)
            }), 500

@app.route('/api/patients/<patient_id>', methods=['GET', 'PUT', 'DELETE'])
def manage_patient(patient_id):
    """Manage individual patient - GET to view, PUT to update, DELETE to remove"""
    if request.method == 'GET':
        try:
            patient_details = practice_db.get_patient_details(patient_id)
            if not patient_details:
                return jsonify({
                    'success': False,
                    'error': 'Patient not found'
                }), 404
            
            return jsonify({
                'success': True,
                **patient_details
            })
        except Exception as e:
            return jsonify({
                'success': False,
                'error': str(e)
            }), 500
    
    elif request.method == 'PUT':
        try:
            patient_data = request.get_json()
            
            # Validate required fields
            if not patient_data.get('first_name') or not patient_data.get('last_name'):
                return jsonify({
                    'success': False,
                    'error': 'First name and last name are required'
                }), 400
            
            success = practice_db.update_patient(patient_id, **patient_data)
            if success:
                return jsonify({
                    'success': True,
                    'message': 'Patient updated successfully'
                })
            else:
                return jsonify({
                    'success': False,
                    'error': 'Patient not found'
                }), 404
        except Exception as e:
            return jsonify({
                'success': False,
                'error': str(e)
            }), 500

@app.route('/api/appointments', methods=['GET', 'POST'])
def manage_appointments():
    """Manage appointments - GET to retrieve, POST to create"""
    if request.method == 'GET':
        week_start = request.args.get('week_start')
        patient_id = request.args.get('patient_id')
        
        try:
            appointments = practice_db.get_appointments(week_start=week_start, patient_id=patient_id)
            return jsonify({
                'success': True,
                'appointments': appointments,
                'count': len(appointments)
            })
        except Exception as e:
            return jsonify({
                'success': False,
                'error': str(e)
            }), 500
    
    elif request.method == 'POST':
        try:
            appointment_data = request.get_json()
            
            # Validate required fields
            required_fields = ['patient_id', 'appointment_date', 'appointment_time']
            for field in required_fields:
                if not appointment_data.get(field):
                    return jsonify({
                        'success': False,
                        'error': f'{field} is required'
                    }), 400
            
            appointment_id = practice_db.create_appointment(**appointment_data)
            return jsonify({
                'success': True,
                'appointment_id': appointment_id,
                'message': 'Appointment created successfully'
            })
        except Exception as e:
            return jsonify({
                'success': False,
                'error': str(e)
            }), 500

@app.route('/api/appointments/<appointment_id>', methods=['PUT', 'DELETE'])
def manage_appointment(appointment_id):
    """Manage specific appointment"""
    if request.method == 'PUT':
        try:
            data = request.get_json()
            status = data.get('status')
            success = practice_db.update_appointment_status(appointment_id, status)
            if success:
                return jsonify({
                    'success': True,
                    'message': 'Rendez-vous mis à jour avec succès'
                })
            else:
                return jsonify({
                    'success': False,
                    'error': 'Rendez-vous non trouvé'
                }), 404
        except Exception as e:
            return jsonify({
                'success': False,
                'error': str(e)
            }), 500
    
    elif request.method == 'DELETE':
        try:
            success = practice_db.delete_appointment(appointment_id)
            if success:
                return jsonify({
                    'success': True,
                    'message': 'Rendez-vous supprimé avec succès'
                })
            else:
                return jsonify({
                    'success': False,
                    'error': 'Rendez-vous non trouvé'
                }), 404
        except Exception as e:
            return jsonify({
                'success': False,
                'error': str(e)
            }), 500

@app.route('/api/schedule-treatment', methods=['POST'])
def schedule_treatment():
    """Enhanced intelligent treatment scheduling"""
    try:
        data = request.json
        patient_id = data.get('patient_id')
        treatment_plan = data.get('treatment_plan', {})
        start_date = data.get('start_date')
        preferred_time = data.get('preferred_time', '09:00')
        use_intelligent_scheduling = data.get('use_intelligent_scheduling', True)
        
        if not patient_id or not treatment_plan or not start_date:
            return jsonify({
                'success': False,
                'error': 'Données manquantes: patient_id, treatment_plan et start_date requis'
            })
        
        treatment_sequence = treatment_plan.get('treatment_sequence', [])
        if not treatment_sequence:
            return jsonify({
                'success': False,
                'error': 'Aucune séquence de traitement trouvée'
            })
        
        # Validate patient exists
        patient = practice_db.get_patient(patient_id)
        if not patient:
            return jsonify({
                'success': False,
                'error': 'Patient non trouvé'
            })
        
        # Use intelligent scheduling if enabled
        if use_intelligent_scheduling:
            print("🧠 Using intelligent scheduling...")
            
            # Generate intelligent schedule
            intelligent_result = intelligent_scheduler.generate_intelligent_schedule(
                patient_id, treatment_sequence, start_date
            )
            
            # Create appointments from intelligent schedule
            appointments = []
            appointment_ids = []
            
            for apt_data in intelligent_result['appointments']:
                appointment_info = {
                    'patient_id': patient_id,
                    'appointment_date': apt_data['date'].strftime('%Y-%m-%d'),
                    'appointment_time': apt_data['time'],
                    'duration_minutes': apt_data['duration_minutes'],
                    'treatment_type': apt_data['treatment'].get('traitement', 'Traitement dentaire'),
                    'doctor': apt_data['treatment'].get('dr', 'Dr.'),
                    'notes': f"{apt_data['treatment'].get('remarque', '')} | {apt_data['reasoning']}",
                    'status': 'scheduled'
                }
                
                # Save appointment to database
                appointment_id = practice_db.create_appointment(**appointment_info)
                appointment_ids.append(appointment_id)
                
                appointments.append({
                    'id': appointment_id,
                    'date': apt_data['date'].strftime('%d/%m/%Y'),
                    'time': apt_data['time'],
                    'treatment': apt_data['treatment'].get('traitement', 'Traitement dentaire'),
                    'duration': f"{apt_data['duration_minutes']} min",
                    'doctor': apt_data['treatment'].get('dr', 'Dr.'),
                    'reasoning': apt_data['reasoning'],
                    'classification': apt_data['classification']['category']
                })
            
            # Save treatment plan to database
            treatment_plan_data = {
                'patient_id': patient_id,
                'consultation_text': treatment_plan.get('consultation_text', ''),
                'treatment_sequence': treatment_sequence,
                'created_at': datetime.now().isoformat(),
                'status': 'scheduled'
            }
            
            plan_id = practice_db.create_treatment_plan(**treatment_plan_data)
            
            return jsonify({
                'success': True,
                'message': f'Traitement programmé intelligemment pour {patient["first_name"]} {patient["last_name"]}',
                'appointments': appointments,
                'treatment_plan_id': plan_id,
                'patient_name': f'{patient["first_name"]} {patient["last_name"]}',
                'intelligent_scheduling': True,
                'scheduling_summary': intelligent_result['scheduling_summary'],
                'ai_analysis': intelligent_result['llm_analysis'],
                'total_duration_minutes': intelligent_result['total_duration']
            })
        
        else:
            # Use original basic scheduling logic
            print("📅 Using basic scheduling...")
            
            # Parse start date
            start_datetime = datetime.strptime(start_date, '%Y-%m-%d')
            
            # Calculate appointment dates based on delays
            appointments = []
            current_date = start_datetime
            
            for i, step in enumerate(treatment_sequence):
                # Extract duration in minutes
                duration_str = step.get('duree', '60 min')
                duration_minutes = 60  # default
                
                if 'min' in duration_str:
                    try:
                        duration_minutes = int(duration_str.split('min')[0].strip())
                    except:
                        duration_minutes = 60
                
                # Calculate appointment date
                if i == 0:
                    # First appointment uses start date
                    appointment_date = current_date
                else:
                    # Subsequent appointments based on delay
                    delay_str = step.get('delai', '1 semaine')
                    delay_days = 7  # default to 1 week
                    
                    if 'semaine' in delay_str:
                        try:
                            weeks = int(delay_str.split('semaine')[0].strip())
                            delay_days = weeks * 7
                        except:
                            delay_days = 7
                    elif 'jour' in delay_str:
                        try:
                            delay_days = int(delay_str.split('jour')[0].strip())
                        except:
                            delay_days = 7
                    elif 'mois' in delay_str:
                        try:
                            months = int(delay_str.split('mois')[0].strip())
                            delay_days = months * 30
                        except:
                            delay_days = 30
                    
                    appointment_date = current_date + timedelta(days=delay_days)
                
                # Skip weekends
                while appointment_date.weekday() >= 5:  # Saturday = 5, Sunday = 6
                    appointment_date += timedelta(days=1)
                
                # Find available time slot
                appointment_time = preferred_time
                available_slots = practice_db.get_available_slots(
                    appointment_date.strftime('%Y-%m-%d'), 
                    duration_minutes
                )
                
                if available_slots:
                    # Try to find preferred time or closest available
                    if preferred_time in available_slots:
                        appointment_time = preferred_time
                    else:
                        # Find closest available slot to preferred time
                        preferred_minutes = int(preferred_time.split(':')[0]) * 60 + int(preferred_time.split(':')[1])
                        closest_slot = min(available_slots, key=lambda slot: 
                            abs((int(slot.split(':')[0]) * 60 + int(slot.split(':')[1])) - preferred_minutes))
                        appointment_time = closest_slot
                else:
                    # No available slots, try next day
                    max_attempts = 7  # Try up to 7 days
                    attempts = 0
                    while not available_slots and attempts < max_attempts:
                        appointment_date += timedelta(days=1)
                        # Skip weekends
                        while appointment_date.weekday() >= 5:
                            appointment_date += timedelta(days=1)
                        available_slots = practice_db.get_available_slots(
                            appointment_date.strftime('%Y-%m-%d'), 
                            duration_minutes
                        )
                        attempts += 1
                    
                    if available_slots:
                        appointment_time = available_slots[0]  # Take first available
                    else:
                        # Force schedule with warning
                        appointment_time = preferred_time
                
                # Create appointment
                appointment_data = {
                    'patient_id': patient_id,
                    'appointment_date': appointment_date.strftime('%Y-%m-%d'),
                    'appointment_time': appointment_time,
                    'duration_minutes': duration_minutes,
                    'treatment_type': step.get('traitement', 'Traitement dentaire'),
                    'doctor': step.get('dr', 'Dr.'),
                    'notes': step.get('remarque', ''),
                    'status': 'scheduled'
                }
                
                # Save appointment to database
                appointment_id = practice_db.create_appointment(**appointment_data)
                
                if appointment_id:
                    appointments.append({
                        'id': appointment_id,
                        'date': appointment_date.strftime('%d/%m/%Y'),
                        'time': appointment_time,
                        'treatment': step.get('traitement', 'Traitement dentaire'),
                        'duration': f"{duration_minutes} min",
                        'doctor': step.get('dr', 'Dr.'),
                        'reasoning': 'Programmation standard',
                        'classification': 'basic'
                    })
                
                # Update current date for next appointment
                current_date = appointment_date
            
            # Save treatment plan to database
            treatment_plan_data = {
                'patient_id': patient_id,
                'consultation_text': treatment_plan.get('consultation_text', ''),
                'treatment_sequence': treatment_sequence,
                'created_at': datetime.now().isoformat(),
                'status': 'scheduled'
            }
            
            plan_id = practice_db.create_treatment_plan(**treatment_plan_data)
            
            return jsonify({
                'success': True,
                'message': f'Traitement programmé avec succès pour {patient["first_name"]} {patient["last_name"]}',
                'appointments': appointments,
                'treatment_plan_id': plan_id,
                'patient_name': f'{patient["first_name"]} {patient["last_name"]}',
                'intelligent_scheduling': False
            })
        
    except Exception as e:
        print(f"❌ Error scheduling treatment: {str(e)}")
        return jsonify({
            'success': False,
            'error': f'Erreur lors de la programmation: {str(e)}'
        })

@app.route('/api/appointments/<appointment_id>/details', methods=['GET'])
def get_appointment_details(appointment_id):
    """Get detailed information about a specific appointment"""
    try:
        conn = sqlite3.connect(practice_db.db_path)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()
        
        cursor.execute('''
            SELECT a.*, p.first_name, p.last_name, p.phone, p.email, p.birth_date,
                   tp.plan_data, tp.consultation_text
            FROM appointments a
            JOIN patients p ON a.patient_id = p.id
            LEFT JOIN treatment_plans tp ON a.treatment_plan_id = tp.id
            WHERE a.id = ?
        ''', (appointment_id,))
        
        appointment = cursor.fetchone()
        conn.close()
        
        if not appointment:
            return jsonify({
                'success': False,
                'error': 'Rendez-vous non trouvé'
            }), 404
        
        appointment_data = dict(appointment)
        
        # Parse treatment plan data if available
        if appointment_data.get('plan_data'):
            try:
                appointment_data['treatment_plan'] = json.loads(appointment_data['plan_data'])
            except:
                appointment_data['treatment_plan'] = None
        
        return jsonify({
            'success': True,
            'appointment': appointment_data
        })
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/api/appointments/<appointment_id>/move', methods=['PUT'])
def move_appointment(appointment_id):
    """Move an appointment to a new date/time"""
    try:
        data = request.json
        new_date = data.get('new_date')
        new_time = data.get('new_time')
        
        if not new_date or not new_time:
            return jsonify({
                'success': False,
                'error': 'Nouvelle date et heure requises'
            }), 400
        
        # Get appointment details first
        conn = sqlite3.connect(practice_db.db_path)
        cursor = conn.cursor()
        
        cursor.execute('SELECT duration_minutes FROM appointments WHERE id = ?', (appointment_id,))
        result = cursor.fetchone()
        
        if not result:
            conn.close()
            return jsonify({
                'success': False,
                'error': 'Rendez-vous non trouvé'
            }), 404
        
        duration_minutes = result[0]
        
        # Check if new slot is available
        available_slots = practice_db.get_available_slots(new_date, duration_minutes)
        
        if new_time not in available_slots:
            # Check if there's a conflict
            cursor.execute('''
                SELECT COUNT(*) FROM appointments 
                WHERE appointment_date = ? AND appointment_time = ? 
                AND id != ? AND status != 'cancelled'
            ''', (new_date, new_time, appointment_id))
            
            conflict_count = cursor.fetchone()[0]
            
            if conflict_count > 0:
                conn.close()
                return jsonify({
                    'success': False,
                    'error': 'Créneau déjà occupé',
                    'available_slots': available_slots
                }), 409
        
        # Update appointment
        cursor.execute('''
            UPDATE appointments 
            SET appointment_date = ?, appointment_time = ?, updated_at = CURRENT_TIMESTAMP
            WHERE id = ?
        ''', (new_date, new_time, appointment_id))
        
        conn.commit()
        conn.close()
        
        return jsonify({
            'success': True,
            'message': 'Rendez-vous déplacé avec succès'
        })
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

# === EXISTING ENDPOINTS ===

@app.route('/chat', methods=['POST'])
def chat():
    try:
        data = request.get_json()
        user_message = data.get('message', '').strip()
        tab = data.get('tab', 'dental-brain')
        
        if not user_message:
            return jsonify({'error': 'Message is required'}), 400
        
        print(f"🔍 Processing message for tab: {tab}")
        print(f"📝 User message: {user_message}")
        
        # Get the specialized LLM for this tab
        if tab not in specialized_llms:
            return jsonify({'error': f'Tab {tab} not supported'}), 400
        
        llm = specialized_llms[tab]
        
        # Get AI response using the specialized LLM
        response_data = llm.generate_response(user_message)
        
        # Check if the response was successful
        if not response_data.get('success', False):
            error_message = response_data.get('error', 'Unknown error occurred')
            print(f"❌ LLM error: {error_message}")
            return jsonify({'error': f'AI response error: {error_message}'}), 500
        
        ai_response = response_data.get('response', '')
        references = response_data.get('references', [])
        context_info = response_data.get('context_info', {})
        
        print(f"🤖 AI response length: {len(ai_response)} characters")
        print(f"📚 References found: {len(references)}")
        
        return jsonify({
            'response': ai_response,
            'references': references,
            'context_info': context_info,
            'tab': tab
        })
        
    except Exception as e:
        print(f"❌ Error in chat endpoint: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'error': f'Internal server error: {str(e)}'}), 500

@app.route('/api/schedule-chat', methods=['POST'])
def schedule_chat():
    """Handle schedule chat messages with autonomous AI decision making"""
    try:
        data = request.get_json()
        user_message = data.get('message', '').strip()
        
        if not user_message:
            return jsonify({'error': 'Message is required'}), 400
        
        print(f"🤖 Processing autonomous schedule request: {user_message}")
        
        # Get current schedule context
        schedule_context = get_current_schedule_context()
        
        # Analyze the schedule request
        analysis = analyze_schedule_request(user_message, schedule_context)
        
        # Generate autonomous rescheduling plan
        autonomous_plan = generate_autonomous_rescheduling_plan(analysis, user_message, schedule_context)
        
        # Generate AI response
        ai_response = generate_autonomous_schedule_response(autonomous_plan, user_message)
        
        # Prepare response data
        response_data = {
            'success': True,
            'response': ai_response,
            'autonomous_plan': autonomous_plan,
            'requires_approval': autonomous_plan.get('requires_approval', True),
            'analysis': {
                'analysis': analysis.get('analysis', 'Demande analysée'),
                'detected_dates': analysis.get('detected_dates', []),
                'immediate_actions': analysis.get('immediate_actions', []),
                'proposed_actions': analysis.get('proposed_actions', []),
                'confirmation_needed': analysis.get('confirmation_needed', '')
            }
        }
        
        print(f"🤖 Autonomous schedule plan generated successfully")
        return jsonify(response_data)
        
    except Exception as e:
        print(f"❌ Error in schedule chat endpoint: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({
            'success': False,
            'error': f'Désolé, une erreur s\'est produite: {str(e)}. Veuillez réessayer.',
            'autonomous_plan': None,
            'requires_approval': False,
            'analysis': {
                'analysis': 'Erreur lors de l\'analyse',
                'detected_dates': [],
                'immediate_actions': [],
                'proposed_actions': [],
                'confirmation_needed': 'Une erreur s\'est produite lors de l\'analyse de votre demande.'
            }
        })

def generate_autonomous_rescheduling_plan(analysis, user_message, schedule_context):
    """Generate a complete autonomous rescheduling plan with AI decisions"""
    try:
        print(f"🧠 Generating autonomous rescheduling plan...")
        
        detected_dates = analysis.get('detected_dates', [])
        if not detected_dates:
            return {
                'success': False,
                'message': 'Aucune date détectée dans votre demande.',
                'decisions': [],
                'requires_approval': False
            }
        
        # Get all appointments for the detected dates
        all_affected_appointments = []
        for date in detected_dates:
            appointments = find_appointments_for_date(date)
            all_affected_appointments.extend(appointments)
        
        if not all_affected_appointments:
            return {
                'success': True,
                'message': f'Aucun rendez-vous trouvé pour les dates {", ".join(detected_dates)}.',
                'decisions': [],
                'requires_approval': False
            }
        
        # Use AI to make intelligent rescheduling decisions
        ai_decisions = make_intelligent_rescheduling_decisions(
            all_affected_appointments, 
            detected_dates, 
            user_message,
            schedule_context
        )
        
        # Calculate statistics
        total_appointments = len(all_affected_appointments)
        successful_reschedules = len([d for d in ai_decisions if d.get('success', False)])
        
        plan = {
            'success': True,
            'message': f'Plan de reprogrammation autonome généré pour {total_appointments} rendez-vous.',
            'decisions': ai_decisions,
            'statistics': {
                'total_appointments': total_appointments,
                'successful_reschedules': successful_reschedules,
                'failed_reschedules': total_appointments - successful_reschedules,
                'affected_dates': detected_dates
            },
            'requires_approval': True,
            'execution_ready': True
        }
        
        print(f"🎯 Autonomous plan generated: {successful_reschedules}/{total_appointments} appointments rescheduled")
        return plan
        
    except Exception as e:
        print(f"❌ Error generating autonomous plan: {e}")
        import traceback
        traceback.print_exc()
        return {
            'success': False,
            'message': f'Erreur lors de la génération du plan: {str(e)}',
            'decisions': [],
            'requires_approval': False
        }

def make_intelligent_rescheduling_decisions(appointments, blocked_dates, user_message, schedule_context):
    """Use AI to make intelligent decisions about where to reschedule each appointment"""
    try:
        print(f"🤖 Making intelligent decisions for {len(appointments)} appointments...")
        
        # Get available slots for the next 2 weeks (excluding blocked dates)
        available_slots = get_available_slots_excluding_dates(blocked_dates, days_ahead=14)
        
        # Use LLM to make intelligent decisions
        decisions_prompt = f"""
Tu es un assistant intelligent de planification dentaire. Tu dois prendre des décisions AUTONOMES pour reprogrammer des rendez-vous.

CONTEXTE ACTUEL:
{schedule_context}

DEMANDE DU DENTISTE: {user_message}

DATES BLOQUÉES: {', '.join(blocked_dates)}

RENDEZ-VOUS À REPROGRAMMER:
{json.dumps([{
    'id': apt.get('id'),
    'patient_name': apt.get('patient_name', 'Patient inconnu'),
    'treatment': apt.get('treatment', 'Traitement'),
    'current_date': apt.get('date'),
    'current_time': apt.get('time'),
    'duration_minutes': apt.get('duration_minutes', 60)
} for apt in appointments], indent=2, ensure_ascii=False)}

CRÉNEAUX DISPONIBLES:
{json.dumps(available_slots, indent=2, ensure_ascii=False)}

INSTRUCTIONS:
1. Pour CHAQUE rendez-vous, choisis le MEILLEUR créneau disponible
2. Prends en compte: type de traitement, préférences patient, optimisation du planning
3. Évite les conflits et optimise les journées
4. Fournis une JUSTIFICATION claire pour chaque décision

RÉPONSE REQUISE (JSON):
{{
    "decisions": [
        {{
            "appointment_id": "id_du_rdv",
            "patient_name": "nom_patient",
            "treatment": "type_traitement",
            "current_slot": "date_actuelle à heure_actuelle",
            "new_date": "YYYY-MM-DD",
            "new_time": "HH:MM",
            "reasoning": "Justification de la décision",
            "confidence": 0.95,
            "success": true
        }}
    ],
    "global_strategy": "Stratégie globale utilisée pour ces reprogrammations",
    "optimization_notes": "Notes sur l'optimisation du planning"
}}
"""
        
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "Tu es un assistant de planification dentaire autonome qui prend des décisions intelligentes."},
                {"role": "user", "content": decisions_prompt}
            ],
            temperature=0.3,
            max_tokens=2000
        )
        
        content = response.choices[0].message.content
        print(f"🤖 AI decisions response: {content[:200]}...")
        
        # Parse AI decisions
        if "```json" in content:
            json_start = content.find("```json") + 7
            json_end = content.find("```", json_start)
            json_content = content[json_start:json_end].strip()
        else:
            json_content = content
        
        try:
            ai_response = json.loads(json_content)
            decisions = ai_response.get('decisions', [])
            
            # Validate and enhance decisions
            validated_decisions = []
            for decision in decisions:
                # Verify the new slot is actually available
                new_date = decision.get('new_date')
                new_time = decision.get('new_time')
                
                if is_slot_available(new_date, new_time, available_slots):
                    decision['success'] = True
                    decision['status'] = 'ready_for_execution'
                else:
                    decision['success'] = False
                    decision['status'] = 'slot_not_available'
                    decision['reasoning'] += ' (ATTENTION: Créneau non disponible)'
                
                validated_decisions.append(decision)
            
            # Add global strategy to decisions
            for decision in validated_decisions:
                decision['global_strategy'] = ai_response.get('global_strategy', 'Optimisation automatique')
            
            print(f"✅ Generated {len(validated_decisions)} intelligent decisions")
            return validated_decisions
            
        except json.JSONDecodeError as e:
            print(f"❌ JSON parsing error: {e}")
            # Fallback to simple rescheduling
            return generate_fallback_decisions(appointments, available_slots)
            
    except Exception as e:
        print(f"❌ Error making intelligent decisions: {e}")
        return generate_fallback_decisions(appointments, available_slots)

def get_available_slots_excluding_dates(blocked_dates, days_ahead=14):
    """Get available slots for the next N days, excluding blocked dates"""
    try:
        available_slots = {}
        today = datetime.now()
        
        for i in range(days_ahead):
            date = today + timedelta(days=i)
            date_str = date.strftime('%Y-%m-%d')
            
            # Skip blocked dates and weekends
            if date_str not in blocked_dates and date.weekday() < 5:
                slots = practice_db.get_available_slots(date_str, 60)
                if slots:
                    day_name = date.strftime('%A')
                    french_day_name = {
                        'Monday': 'Lundi',
                        'Tuesday': 'Mardi', 
                        'Wednesday': 'Mercredi',
                        'Thursday': 'Jeudi',
                        'Friday': 'Vendredi'
                    }.get(day_name, day_name)
                    
                    available_slots[date_str] = {
                        'date': date_str,
                        'day_name': f"{french_day_name} {date.strftime('%d/%m')}",
                        'slots': slots
                    }
        
        return available_slots
        
    except Exception as e:
        print(f"❌ Error getting available slots: {e}")
        return {}

def is_slot_available(date, time, available_slots):
    """Check if a specific slot is available"""
    if not date or not time or not available_slots:
        return False
    
    date_slots = available_slots.get(date, {})
    return time in date_slots.get('slots', [])

def generate_fallback_decisions(appointments, available_slots):
    """Generate simple fallback decisions if AI fails"""
    decisions = []
    
    # Simple strategy: assign appointments to first available slots
    slot_list = []
    for date, info in available_slots.items():
        for time in info['slots']:
            slot_list.append({
                'date': date,
                'time': time,
                'day_name': info['day_name']
            })
    
    for i, appointment in enumerate(appointments):
        if i < len(slot_list):
            slot = slot_list[i]
            decisions.append({
                'appointment_id': appointment.get('id'),
                'patient_name': appointment.get('patient_name', 'Patient inconnu'),
                'treatment': appointment.get('treatment', 'Traitement'),
                'current_slot': f"{appointment.get('date')} à {appointment.get('time')}",
                'new_date': slot['date'],
                'new_time': slot['time'],
                'reasoning': f"Reprogrammation automatique vers le premier créneau disponible ({slot['day_name']})",
                'confidence': 0.7,
                'success': True,
                'status': 'ready_for_execution'
            })
        else:
            decisions.append({
                'appointment_id': appointment.get('id'),
                'patient_name': appointment.get('patient_name', 'Patient inconnu'),
                'treatment': appointment.get('treatment', 'Traitement'),
                'current_slot': f"{appointment.get('date')} à {appointment.get('time')}",
                'new_date': None,
                'new_time': None,
                'reasoning': 'Aucun créneau disponible trouvé',
                'confidence': 0.0,
                'success': False,
                'status': 'no_slots_available'
            })
    
    return decisions

def generate_autonomous_schedule_response(autonomous_plan, user_message):
    """Generate a response for the autonomous scheduling plan"""
    try:
        if not autonomous_plan.get('success', False):
            return autonomous_plan.get('message', 'Impossible de générer un plan de reprogrammation.')
        
        decisions = autonomous_plan.get('decisions', [])
        stats = autonomous_plan.get('statistics', {})
        
        if not decisions:
            return "Aucun rendez-vous à reprogrammer trouvé."
        
        # Build response
        response = f"🤖 **Plan de reprogrammation autonome généré**\n\n"
        response += f"J'ai analysé votre demande et préparé un plan complet pour reprogrammer vos {stats.get('total_appointments', 0)} rendez-vous.\n\n"
        
        # Show successful reschedules
        successful = [d for d in decisions if d.get('success', False)]
        if successful:
            response += f"✅ **{len(successful)} rendez-vous reprogrammés avec succès :**\n\n"
            for decision in successful:
                response += f"**{decision.get('patient_name')}** - {decision.get('treatment')}\n"
                response += f"• De : {decision.get('current_slot')}\n"
                response += f"• Vers : {decision.get('new_date')} à {decision.get('new_time')}\n"
                response += f"• Justification : {decision.get('reasoning')}\n"
                response += f"• Confiance : {decision.get('confidence', 0):.0%}\n\n"
        
        # Show failed reschedules
        failed = [d for d in decisions if not d.get('success', False)]
        if failed:
            response += f"⚠️ **{len(failed)} rendez-vous nécessitent attention :**\n\n"
            for decision in failed:
                response += f"**{decision.get('patient_name')}** - {decision.get('treatment')}\n"
                response += f"• Problème : {decision.get('reasoning')}\n\n"
        
        response += f"\n🎯 **Prêt pour exécution** - Cliquez sur 'Approuver' pour appliquer ces modifications à votre planning."
        
        return response
        
    except Exception as e:
        print(f"❌ Error generating autonomous response: {e}")
        return "Plan de reprogrammation généré. Veuillez vérifier les détails."

@app.route('/')
def index():
    """Serve the main chat interface"""
    return render_template('index.html')

@app.route('/modify-treatment-plan', methods=['POST'])
def modify_treatment_plan():
    """Modify treatment plan with natural language"""
    try:
        data = request.get_json()
        modification_request = data.get('modification', '')
        current_plan = data.get('current_plan', {})
        
        if not modification_request:
            return jsonify({'error': 'Modification request is required'}), 400
        
        # Use dental-brain LLM for modifications
        llm = specialized_llms.get('dental-brain')
        if not llm:
            return jsonify({'error': 'Dental brain LLM not available'}), 500
        
        # Create a prompt for modification
        prompt = f"""
        Modifiez le plan de traitement suivant selon cette demande: {modification_request}
        
        Plan actuel: {json.dumps(current_plan, ensure_ascii=False)}
        
        Retournez le plan modifié au format JSON avec la même structure.
        """
        
        response_data = llm.generate_response(prompt)
        
        # Check if the response was successful
        if not response_data.get('success', False):
            error_message = response_data.get('error', 'Unknown error occurred')
            return jsonify({'error': f'AI modification error: {error_message}'}), 500
        
        return jsonify({
            'success': True,
            'modified_plan': response_data.get('response', '')
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/get-reference-details/<reference_id>')
def get_reference_details(reference_id):
    """Get detailed information about a reference"""
    try:
        # Search for the reference in both collections
        reference_details = None
        
        # Try to find in cases collection
        try:
            cases_results = rag_system.cases_collection.get(ids=[reference_id])
            if cases_results['documents'] and len(cases_results['documents']) > 0:
                reference_details = {
                    'id': reference_id,
                    'title': cases_results['metadatas'][0].get('consultation', 'Cas clinique'),
                    'type': 'case',
                    'category': 'Cas clinique',
                    'content': cases_results['documents'][0],
                    'metadata': cases_results['metadatas'][0]
                }
        except Exception as e:
            pass
        
        # Try to find in knowledge collection if not found in cases
        if not reference_details:
            try:
                knowledge_results = rag_system.knowledge_collection.get(ids=[reference_id])
                if knowledge_results['documents'] and len(knowledge_results['documents']) > 0:
                    reference_details = {
                        'id': reference_id,
                        'title': knowledge_results['metadatas'][0].get('title', 'Connaissance dentaire'),
                        'type': knowledge_results['metadatas'][0].get('type', 'knowledge'),
                        'category': knowledge_results['metadatas'][0].get('category', 'Connaissance générale'),
                        'content': knowledge_results['documents'][0],
                        'metadata': knowledge_results['metadatas'][0]
                    }
            except Exception as e:
                pass
        
        if reference_details:
            return jsonify({
                'success': True,
                'reference': reference_details
            })
        else:
            return jsonify({
                'success': False,
                'error': 'Reference not found'
            }), 404
            
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/health')
def health_check():
    """Health check endpoint"""
    try:
        stats = rag_system.get_collection_stats()
        return jsonify({
            'status': 'healthy',
            'rag_enabled': True,
            'openai_configured': bool(api_key),
            'cases_count': stats['cases_count'],
            'knowledge_count': stats['knowledge_count'],
            'total_items': stats['total_items'],
            'timestamp': datetime.now().isoformat()
        })
    except Exception as e:
        return jsonify({
            'status': 'error',
            'error': str(e),
            'timestamp': datetime.now().isoformat()
        }), 500

@app.route('/knowledge')
def get_knowledge_stats():
    """Get knowledge base statistics"""
    try:
        stats = rag_system.get_collection_stats()
        return jsonify({
            'success': True,
            'cases': stats['cases_count'],
            'knowledge': stats['knowledge_count'],
            'total': stats['total_items'],
            'collections': stats['collections']
        })
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e),
            'cases': 0,
            'knowledge': 0,
            'total': 0
        })

@app.route('/search', methods=['POST'])
def search_knowledge():
    """Search endpoint for testing RAG system"""
    try:
        data = request.get_json()
        query = data.get('query', '')
        search_type = data.get('type', 'combined')  # 'cases', 'knowledge', or 'combined'
        
        if not query:
            return jsonify({'success': False, 'error': 'Query is required'})
        
        if search_type == 'cases':
            results = rag_system.search_cases(query, n_results=5)
            return jsonify({
                'success': True,
                'query': query,
                'type': 'cases',
                'results': results,
                'count': len(results)
            })
        elif search_type == 'knowledge':
            results = rag_system.search_knowledge(query, n_results=5)
            return jsonify({
                'success': True,
                'query': query,
                'type': 'knowledge',
                'results': results,
                'count': len(results)
            })
        else:  # combined
            results = rag_system.search_combined(query, case_results=3, knowledge_results=5)
            return jsonify({
                'success': True,
                'query': query,
                'type': 'combined',
                'results': results,
                'count': results['total_results']
            })
            
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/reindex', methods=['POST'])
def reindex_knowledge():
    """Reindex all knowledge"""
    try:
        results = rag_system.reindex_all()
        return jsonify({
            'success': True,
            'message': 'Reindexing completed',
            'results': results
        })
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/static/<path:filename>')
def serve_static(filename):
    """Serve static files"""
    return send_from_directory('static', filename)

@app.route('/debug/context', methods=['POST'])
def debug_context():
    """Debug endpoint to see RAG context for a query"""
    try:
        data = request.get_json()
        query = data.get('query', '')
        
        if not query:
            return jsonify({'success': False, 'error': 'Query is required'})
        
        # Get RAG results
        rag_results = rag_system.search_combined(query, case_results=2, knowledge_results=4)
        
        return jsonify({
            'success': True,
            'query': query,
            'rag_results': rag_results,
            'context_preview': {
                'cases_count': len(rag_results['cases']),
                'knowledge_count': len(rag_results['knowledge']),
                'total_items': rag_results['total_results']
            }
        })
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/debug/prompt-size', methods=['POST'])
def debug_prompt_size():
    """Debug endpoint to see actual prompt sizes for each tab"""
    try:
        data = request.get_json()
        user_message = data.get('message', 'Plan de traitement pour une dent 26 avec carie profonde')
        tab = data.get('tab', 'dental-brain')
        
        if tab not in specialized_llms:
            return jsonify({'error': f'Tab {tab} not supported'}), 400
        
        llm = specialized_llms[tab]
        
        # Get specialized context (same as in generate_response)
        context, references = llm.get_specialized_context(user_message)
        
        # Build complete system prompt
        system_prompt = llm.base_system_prompt
        if context:
            system_prompt += f"\n\n=== CONTEXTE SPÉCIALISÉ ===\n{context}"
        
        # Prepare messages
        messages = [
            {"role": "system", "content": system_prompt},
            *llm.chat_history,
            {"role": "user", "content": user_message}
        ]
        
        # Calculate sizes
        total_tokens = 0
        message_details = []
        
        for msg in messages:
            content = msg['content']
            # Rough token estimation (1 token ≈ 4 characters for French)
            estimated_tokens = len(content) // 4
            total_tokens += estimated_tokens
            
            message_details.append({
                'role': msg['role'],
                'length': len(content),
                'estimated_tokens': estimated_tokens,
                'content_preview': content[:200] + "..." if len(content) > 200 else content
            })
        
        return jsonify({
            'success': True,
            'tab': tab,
            'user_message': user_message,
            'total_estimated_tokens': total_tokens,
            'message_count': len(messages),
            'breakdown': {
                'base_system_prompt_length': len(llm.base_system_prompt),
                'context_length': len(context),
                'references_count': len(references),
                'chat_history_messages': len(llm.chat_history),
                'user_message_length': len(user_message)
            },
            'messages': message_details,
            'context_preview': context[:500] + "..." if len(context) > 500 else context
        })
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/export-treatment-plan', methods=['POST'])
def export_treatment_plan():
    """Export treatment plan as PDF document"""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib import colors
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.units import inch
        from io import BytesIO
        
        data = request.get_json()
        patient_info = data.get('patient_info', {})
        treatment_sequence = data.get('treatment_sequence', [])
        consultation_text = data.get('consultation_text', '')
        
        if not treatment_sequence:
            return jsonify({'success': False, 'error': 'Treatment sequence is required'}), 400
        
        # Create PDF in memory
        buffer = BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=1*inch, bottomMargin=1*inch)
        
        # Get styles
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=30,
            alignment=1,  # Center alignment
            textColor=colors.HexColor('#2c5aa0')
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=14,
            spaceAfter=12,
            textColor=colors.HexColor('#2c5aa0')
        )
        
        normal_style = styles['Normal']
        
        # Build document content
        content = []
        
        # Title
        content.append(Paragraph("Plan de Traitement Dentaire", title_style))
        content.append(Spacer(1, 20))
        
        # Document info
        current_date = datetime.now().strftime("%d/%m/%Y")
        
        info_data = [
            ['Date:', current_date],
            ['Généré par:', 'Dental AI Assistant']
        ]
        
        # Add patient info if available
        if patient_info:
            for key, value in patient_info.items():
                display_key = key.replace('_', ' ').title()
                info_data.append([f'{display_key}:', str(value)])
        
        info_table = Table(info_data, colWidths=[1.5*inch, 4*inch])
        info_table.setStyle(TableStyle([
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ]))
        
        content.append(info_table)
        content.append(Spacer(1, 20))
        
        # Add consultation text if available
        if consultation_text:
            content.append(Paragraph("Consultation", heading_style))
            content.append(Spacer(1, 8))
            
            # Format consultation text
            consultation_paragraph = Paragraph(consultation_text, normal_style)
            content.append(consultation_paragraph)
            content.append(Spacer(1, 20))
        
        # Treatment sequence table
        content.append(Paragraph("Séquence de Traitement", heading_style))
        content.append(Spacer(1, 12))
        
        if treatment_sequence:
            # Table headers
            table_data = [['RDV', 'Traitement', 'Durée', 'Délai', 'Praticien', 'Remarques']]
            
            # Table rows
            for treatment in treatment_sequence:
                row = [
                    str(treatment.get('rdv', treatment.get('step', ''))),
                    treatment.get('traitement', treatment.get('treatment', '')),
                    treatment.get('duree', treatment.get('duration', '')),
                    treatment.get('delai', treatment.get('delay', '')),
                    treatment.get('dr', treatment.get('doctor', '')),
                    treatment.get('remarque', treatment.get('remarks', ''))
                ]
                table_data.append(row)
            
            # Create table
            treatment_table = Table(table_data, colWidths=[0.6*inch, 2.2*inch, 1*inch, 1*inch, 1*inch, 1.7*inch])
            treatment_table.setStyle(TableStyle([
                # Header styling
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2c5aa0')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                
                # Body styling
                ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                ('FONTSIZE', (0, 1), (-1, -1), 9),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.beige, colors.white]),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ('LEFTPADDING', (0, 0), (-1, -1), 6),
                ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                ('TOPPADDING', (0, 0), (-1, -1), 8),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
            ]))
            
            content.append(treatment_table)
            content.append(Spacer(1, 30))
            
            # Summary
            content.append(Paragraph("Résumé", heading_style))
            summary_text = f"""
            <b>Nombre total de rendez-vous:</b> {len(treatment_sequence)}<br/>
            <b>Plan généré le:</b> {current_date}<br/>
            <b>Assistant:</b> Dental AI - Intelligence Artificielle Dentaire
            """
            content.append(Paragraph(summary_text, normal_style))
        
        # Build PDF
        doc.build(content)
        
        # Get PDF data
        pdf_data = buffer.getvalue()
        buffer.close()
        
        # Create response
        response = app.response_class(
            pdf_data,
            mimetype='application/pdf',
            headers={
                'Content-Disposition': 'attachment; filename=plan-traitement.pdf',
                'Content-Length': len(pdf_data)
            }
        )
        
        return response
        
    except ImportError:
        return jsonify({
            'success': False,
            'error': 'PDF generation library not installed. Please install reportlab.'
        }), 500
    except Exception as e:
        print(f"❌ Error generating PDF: {e}")
        return jsonify({
            'success': False,
            'error': f'Erreur lors de la génération du PDF: {str(e)}'
        }), 500

# === FINANCIAL MANAGEMENT ENDPOINTS ===

@app.route('/api/pricing', methods=['GET'])
def get_pricing():
    """Get dental pricing data"""
    try:
        search_term = request.args.get('search', '')
        pricing_data = practice_db.get_pricing_data(search_term if search_term else None)
        
        return jsonify({
            'success': True,
            'pricing': pricing_data,
            'count': len(pricing_data)
        })
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/api/invoices', methods=['GET', 'POST'])
def manage_invoices():
    """Manage invoices - GET to retrieve, POST to create"""
    if request.method == 'GET':
        try:
            patient_id = request.args.get('patient_id')
            status = request.args.get('status')
            
            invoices = practice_db.get_invoices(patient_id=patient_id, status=status)
            return jsonify({
                'success': True,
                'invoices': invoices,
                'count': len(invoices)
            })
        except Exception as e:
            return jsonify({
                'success': False,
                'error': str(e)
            }), 500
    
    elif request.method == 'POST':
        try:
            data = request.get_json()
            patient_id = data.get('patient_id')
            treatment_items = data.get('treatment_items', [])
            invoice_date = data.get('invoice_date')
            due_date = data.get('due_date')
            
            if not patient_id or not treatment_items:
                return jsonify({
                    'success': False,
                    'error': 'Patient ID and treatment items are required'
                }), 400
            
            invoice_id = practice_db.create_invoice(
                patient_id=patient_id,
                treatment_items=treatment_items,
                invoice_date=invoice_date,
                due_date=due_date
            )
            
            if invoice_id:
                return jsonify({
                    'success': True,
                    'invoice_id': invoice_id,
                    'message': 'Facture créée avec succès'
                })
            else:
                return jsonify({
                    'success': False,
                    'error': 'Erreur lors de la création de la facture'
                }), 500
                
        except Exception as e:
            return jsonify({
                'success': False,
                'error': str(e)
            }), 500

@app.route('/api/payments', methods=['POST'])
def add_payment():
    """Add a payment to an invoice"""
    try:
        data = request.get_json()
        invoice_id = data.get('invoice_id')
        amount = data.get('amount')
        payment_date = data.get('payment_date')
        payment_method = data.get('payment_method', 'cash')
        reference_number = data.get('reference_number')
        
        if not invoice_id or not amount:
            return jsonify({
                'success': False,
                'error': 'Invoice ID and amount are required'
            }), 400
        
        payment_id = practice_db.add_payment(
            invoice_id=invoice_id,
            amount=amount,
            payment_date=payment_date,
            payment_method=payment_method,
            reference_number=reference_number
        )
        
        if payment_id:
            return jsonify({
                'success': True,
                'payment_id': payment_id,
                'message': 'Paiement enregistré avec succès'
            })
        else:
            return jsonify({
                'success': False,
                'error': 'Erreur lors de l\'enregistrement du paiement'
            }), 500
            
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/api/financial-dashboard', methods=['GET'])
def get_financial_dashboard():
    """Get financial dashboard data for analytics"""
    try:
        dashboard_data = practice_db.get_financial_dashboard_data()
        
        # Add some calculated metrics
        current_month_revenue = 0
        previous_month_revenue = 0
        
        if dashboard_data.get('monthly_revenue'):
            monthly_data = dashboard_data['monthly_revenue']
            if len(monthly_data) >= 1:
                current_month_revenue = monthly_data[-1]['revenue']
            if len(monthly_data) >= 2:
                previous_month_revenue = monthly_data[-2]['revenue']
        
        # Calculate growth
        growth_rate = 0
        if previous_month_revenue > 0:
            growth_rate = ((current_month_revenue - previous_month_revenue) / previous_month_revenue) * 100
        
        # Calculate total pending payments
        total_pending = 0
        total_paid = 0
        for status_data in dashboard_data.get('payment_status', []):
            if status_data['status'] == 'pending':
                total_pending = status_data['amount']
            elif status_data['status'] == 'paid':
                total_paid = status_data['amount']
        
        dashboard_data['summary'] = {
            'current_month_revenue': current_month_revenue,
            'previous_month_revenue': previous_month_revenue,
            'growth_rate': growth_rate,
            'total_pending': total_pending,
            'total_paid': total_paid,
            'collection_rate': (total_paid / (total_paid + total_pending)) * 100 if (total_paid + total_pending) > 0 else 0
        }
        
        return jsonify({
            'success': True,
            'dashboard': dashboard_data
        })
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

# === DEVIS (ESTIMATES) ENDPOINTS ===

@app.route('/api/devis', methods=['GET', 'POST'])
def manage_devis():
    """Manage devis - GET to retrieve, POST to create"""
    if request.method == 'GET':
        try:
            patient_id = request.args.get('patient_id')
            status = request.args.get('status')
            devis_id = request.args.get('devis_id')
            
            devis_list = practice_db.get_devis(patient_id=patient_id, status=status, devis_id=devis_id)
            return jsonify({
                'success': True,
                'devis': devis_list,
                'count': len(devis_list)
            })
        except Exception as e:
            return jsonify({
                'success': False,
                'error': str(e)
            }), 500
    
    elif request.method == 'POST':
        try:
            data = request.get_json()
            patient_id = data.get('patient_id')
            treatment_plan_id = data.get('treatment_plan_id')
            devis_items = data.get('devis_items', [])
            valid_days = data.get('valid_days', 30)
            
            if not patient_id or not devis_items:
                return jsonify({
                    'success': False,
                    'error': 'Patient ID and devis items are required'
                }), 400
            
            devis_id = practice_db.create_devis(
                patient_id=patient_id,
                treatment_plan_id=treatment_plan_id,
                devis_items=devis_items,
                valid_days=valid_days
            )
            
            if devis_id:
                return jsonify({
                    'success': True,
                    'devis_id': devis_id,
                    'message': 'Devis créé avec succès'
                })
            else:
                return jsonify({
                    'success': False,
                    'error': 'Erreur lors de la création du devis'
                }), 500
                
        except Exception as e:
            return jsonify({
                'success': False,
                'error': str(e)
            }), 500

@app.route('/api/devis/<devis_id>/approve', methods=['POST'])
def approve_devis(devis_id):
    """Approve a devis"""
    try:
        success = practice_db.approve_devis(devis_id)
        
        if success:
            return jsonify({
                'success': True,
                'message': 'Devis approuvé avec succès'
            })
        else:
            return jsonify({
                'success': False,
                'error': 'Erreur lors de l\'approbation du devis'
            }), 500
            
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/api/devis/<devis_id>/reject', methods=['POST'])
def reject_devis(devis_id):
    """Reject a devis"""
    try:
        data = request.get_json()
        reason = data.get('reason', '')
        
        success = practice_db.reject_devis(devis_id, reason)
        
        if success:
            return jsonify({
                'success': True,
                'message': 'Devis rejeté'
            })
        else:
            return jsonify({
                'success': False,
                'error': 'Erreur lors du rejet du devis'
            }), 500
            
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/api/devis/<devis_id>/create-invoice', methods=['POST'])
def create_invoice_from_devis(devis_id):
    """Create an invoice from an approved devis"""
    try:
        data = request.get_json()
        selected_items = data.get('selected_items')  # Optional: specific items to invoice
        
        invoice_id = practice_db.create_invoice_from_devis(devis_id, selected_items)
        
        if invoice_id:
            return jsonify({
                'success': True,
                'invoice_id': invoice_id,
                'message': 'Facture créée à partir du devis'
            })
        else:
            return jsonify({
                'success': False,
                'error': 'Erreur lors de la création de la facture'
            }), 500
            
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/api/generate-devis-from-treatment', methods=['POST'])
def generate_devis_from_treatment():
    """Generate devis from treatment plan"""
    try:
        data = request.get_json()
        patient_id = data.get('patient_id')
        treatment_plan_id = data.get('treatment_plan_id')
        treatment_plan = data.get('treatment_plan')
        
        if not patient_id or not treatment_plan:
            return jsonify({
                'success': False,
                'error': 'Patient ID and treatment plan are required'
            }), 400
        
        # Get pricing data
        pricing_data = practice_db.get_pricing_data()
        pricing_map = {item['treatment_name'].lower(): item for item in pricing_data}
        
        # Convert treatment plan to devis items
        devis_items = []
        for step in treatment_plan.get('treatment_sequence', []):
            treatment_name = step.get('traitement', '').lower()
            
            # Try to find matching pricing
            matched_pricing = None
            for pricing_key, pricing_item in pricing_map.items():
                if any(keyword in treatment_name for keyword in pricing_key.split()):
                    matched_pricing = pricing_item
                    break
            
            if matched_pricing:
                devis_items.append({
                    'tarmed_code': matched_pricing['tarmed_code'],
                    'treatment_name': matched_pricing['treatment_name'],
                    'quantity': 1,
                    'unit_price': matched_pricing['base_price_chf'],
                    'lamal_covered': matched_pricing['lamal_covered'],
                    'lamal_percentage': matched_pricing['lamal_percentage'],
                    'discount_percentage': 0,
                    'discount_amount_chf': 0,
                    'notes': step.get('remarque', '')
                })
            else:
                # Default pricing if no match found
                devis_items.append({
                    'tarmed_code': '00.0000',
                    'treatment_name': step.get('traitement', 'Traitement non spécifié'),
                    'quantity': 1,
                    'unit_price': 200.0,  # Default price
                    'lamal_covered': False,
                    'lamal_percentage': 0,
                    'discount_percentage': 0,
                    'discount_amount_chf': 0,
                    'notes': step.get('remarque', '')
                })
        
        # Create devis
        devis_id = practice_db.create_devis(
            patient_id=patient_id,
            treatment_plan_id=treatment_plan_id,
            devis_items=devis_items
        )
        
        if devis_id:
            return jsonify({
                'success': True,
                'devis_id': devis_id,
                'devis_items': devis_items,
                'message': 'Devis généré avec succès'
            })
        else:
            return jsonify({
                'success': False,
                'error': 'Erreur lors de la génération du devis'
            }), 500
            
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

# === PAYMENT PLANS ENDPOINTS ===

@app.route('/api/payment-plans', methods=['GET', 'POST'])
def manage_payment_plans():
    """Manage payment plans"""
    if request.method == 'GET':
        try:
            invoice_id = request.args.get('invoice_id')
            payment_plans = practice_db.get_payment_plans(invoice_id=invoice_id)
            
            return jsonify({
                'success': True,
                'payment_plans': payment_plans,
                'count': len(payment_plans)
            })
        except Exception as e:
            return jsonify({
                'success': False,
                'error': str(e)
            }), 500
    
    elif request.method == 'POST':
        try:
            data = request.get_json()
            invoice_id = data.get('invoice_id')
            plan_name = data.get('plan_name')
            number_of_payments = data.get('number_of_payments')
            frequency = data.get('frequency', 'monthly')
            first_payment_date = data.get('first_payment_date')
            
            if not invoice_id or not plan_name or not number_of_payments:
                return jsonify({
                    'success': False,
                    'error': 'Invoice ID, plan name, and number of payments are required'
                }), 400
            
            plan_id = practice_db.create_payment_plan(
                invoice_id=invoice_id,
                plan_name=plan_name,
                number_of_payments=number_of_payments,
                frequency=frequency,
                first_payment_date=first_payment_date
            )
            
            if plan_id:
                return jsonify({
                    'success': True,
                    'plan_id': plan_id,
                    'message': 'Plan de paiement créé avec succès'
                })
            else:
                return jsonify({
                    'success': False,
                    'error': 'Erreur lors de la création du plan de paiement'
                }), 500
                
        except Exception as e:
            return jsonify({
                'success': False,
                'error': str(e)
            }), 500

@app.route('/api/revenue-forecast', methods=['GET'])
def get_revenue_forecast():
    """Get revenue forecast"""
    try:
        months_ahead = request.args.get('months', 12, type=int)
        forecast = practice_db.get_revenue_forecast(months_ahead=months_ahead)
        
        return jsonify({
            'success': True,
            'forecast': forecast
        })
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/api/generate-treatment-invoice', methods=['POST'])
def generate_treatment_invoice():
    """Generate invoice from treatment plan"""
    try:
        data = request.get_json()
        patient_id = data.get('patient_id')
        treatment_plan = data.get('treatment_plan')
        treatment_plan_id = data.get('treatment_plan_id')
        
        if not patient_id or not treatment_plan:
            return jsonify({
                'success': False,
                'error': 'Patient ID and treatment plan are required'
            }), 400
        
        # Get pricing data
        pricing_data = practice_db.get_pricing_data()
        pricing_map = {item['treatment_name'].lower(): item for item in pricing_data}
        
        # Convert treatment plan to invoice items
        invoice_items = []
        for step in treatment_plan.get('treatment_sequence', []):
            treatment_name = step.get('traitement', '').lower()
            
            # Try to find matching pricing
            matched_pricing = None
            for pricing_key, pricing_item in pricing_map.items():
                if any(keyword in treatment_name for keyword in pricing_key.split()):
                    matched_pricing = pricing_item
                    break
            
            if matched_pricing:
                invoice_items.append({
                    'tarmed_code': matched_pricing['tarmed_code'],
                    'treatment_name': matched_pricing['treatment_name'],
                    'quantity': 1,
                    'unit_price': matched_pricing['base_price_chf'],
                    'lamal_covered': matched_pricing['lamal_covered'],
                    'lamal_percentage': matched_pricing['lamal_percentage']
                })
            else:
                # Default pricing if no match found
                invoice_items.append({
                    'tarmed_code': '00.0000',
                    'treatment_name': step.get('traitement', 'Traitement non spécifié'),
                    'quantity': 1,
                    'unit_price': 200.0,  # Default price
                    'lamal_covered': False,
                    'lamal_percentage': 0
                })
        
        # Create invoice
        invoice_id = practice_db.create_invoice(
            patient_id=patient_id,
            treatment_items=invoice_items,
            treatment_plan_id=treatment_plan_id
        )
        
        if invoice_id:
            return jsonify({
                'success': True,
                'invoice_id': invoice_id,
                'invoice_items': invoice_items,
                'message': 'Facture générée avec succès'
            })
        else:
            return jsonify({
                'success': False,
                'error': 'Erreur lors de la génération de la facture'
            }), 500
            
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/api/devis/<devis_id>/download', methods=['GET'])
def download_devis_pdf(devis_id):
    """Download devis as PDF"""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib import colors
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.units import inch
        from io import BytesIO
        
        # Get devis data
        devis_list = practice_db.get_devis(devis_id=devis_id)
        if not devis_list:
            return jsonify({'success': False, 'error': 'Devis non trouvé'}), 404
        
        devis = devis_list[0]
        
        # Create PDF in memory
        buffer = BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=1*inch, bottomMargin=1*inch)
        
        # Get styles
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=30,
            alignment=1,  # Center alignment
            textColor=colors.HexColor('#2c5aa0')
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=14,
            spaceAfter=12,
            textColor=colors.HexColor('#2c5aa0')
        )
        
        normal_style = styles['Normal']
        
        # Build document content
        content = []
        
        # Title
        content.append(Paragraph("DEVIS DENTAIRE", title_style))
        content.append(Spacer(1, 20))
        
        # Devis info
        devis_info = f"""
        <b>Numéro de devis:</b> {devis['devis_number']}<br/>
        <b>Date:</b> {datetime.strptime(devis['devis_date'], '%Y-%m-%d').strftime('%d/%m/%Y')}<br/>
        <b>Valide jusqu'au:</b> {datetime.strptime(devis['valid_until'], '%Y-%m-%d').strftime('%d/%m/%Y')}<br/>
        <b>Patient:</b> {devis['patient_name']}<br/>
        <b>Statut:</b> {devis['status'].upper()}
        """
        content.append(Paragraph(devis_info, normal_style))
        content.append(Spacer(1, 20))
        
        # Treatments table
        content.append(Paragraph("Détail des traitements", heading_style))
        
        # Get devis items
        devis_items = practice_db.get_devis_items(devis_id)
        
        if devis_items:
            # Table headers
            table_data = [['Code TARMED', 'Traitement', 'Qté', 'Prix unitaire', 'Total', 'LAMal']]
            
            # Table rows
            for item in devis_items:
                row = [
                    item['tarmed_code'],
                    item['treatment_name'],
                    str(item['quantity']),
                    f"{item['unit_price_chf']:.2f} CHF",
                    f"{item['final_price_chf']:.2f} CHF",
                    'Oui' if item['lamal_covered'] else 'Non'
                ]
                table_data.append(row)
            
            # Create table
            treatments_table = Table(table_data, colWidths=[1*inch, 2.5*inch, 0.5*inch, 1*inch, 1*inch, 0.8*inch])
            treatments_table.setStyle(TableStyle([
                # Header styling
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2c5aa0')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                
                # Body styling
                ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                ('FONTSIZE', (0, 1), (-1, -1), 9),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.beige, colors.white]),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ('LEFTPADDING', (0, 0), (-1, -1), 6),
                ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                ('TOPPADDING', (0, 0), (-1, -1), 8),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
            ]))
            
            content.append(treatments_table)
            content.append(Spacer(1, 20))
        
        # Totals
        content.append(Paragraph("Récapitulatif", heading_style))
        totals_text = f"""
        <b>Total:</b> {devis['total_amount_chf']:.2f} CHF<br/>
        <b>Prise en charge LAMal:</b> {devis['lamal_amount_chf']:.2f} CHF<br/>
        <b>Assurance complémentaire:</b> {devis['insurance_amount_chf']:.2f} CHF<br/>
        <b>À payer par le patient:</b> {devis['patient_amount_chf']:.2f} CHF
        """
        content.append(Paragraph(totals_text, normal_style))
        content.append(Spacer(1, 30))
        
        # Footer
        footer_text = f"""
        <i>Devis généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')}</i><br/>
        <i>Assistant: Dental AI - Intelligence Artificielle Dentaire</i>
        """
        content.append(Paragraph(footer_text, normal_style))
        
        # Build PDF
        doc.build(content)
        
        # Get PDF data
        pdf_data = buffer.getvalue()
        buffer.close()
        
        # Create response
        response = app.response_class(
            pdf_data,
            mimetype='application/pdf',
            headers={
                'Content-Disposition': f'attachment; filename=devis-{devis["devis_number"]}.pdf',
                'Content-Length': len(pdf_data)
            }
        )
        
        return response
        
    except ImportError:
        return jsonify({
            'success': False,
            'error': 'PDF generation library not installed. Please install reportlab.'
        }), 500
    except Exception as e:
        print(f"❌ Error generating devis PDF: {e}")
        return jsonify({
            'success': False,
            'error': f'Erreur lors de la génération du PDF: {str(e)}'
        }), 500

@app.route('/api/invoices/<invoice_id>/download', methods=['GET'])
def download_invoice_pdf(invoice_id):
    """Download invoice as PDF"""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib import colors
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.units import inch
        from io import BytesIO
        
        # Get invoice data
        invoices = practice_db.get_invoices(invoice_id=invoice_id)
        if not invoices:
            return jsonify({'success': False, 'error': 'Facture non trouvée'}), 404
        
        invoice = invoices[0]
        
        # Create PDF in memory
        buffer = BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=1*inch, bottomMargin=1*inch)
        
        # Get styles
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=30,
            alignment=1,  # Center alignment
            textColor=colors.HexColor('#2c5aa0')
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=14,
            spaceAfter=12,
            textColor=colors.HexColor('#2c5aa0')
        )
        
        normal_style = styles['Normal']
        
        # Build document content
        content = []
        
        # Title
        content.append(Paragraph("FACTURE", title_style))
        content.append(Spacer(1, 20))
        
        # Invoice info
        invoice_info = f"""
        <b>Numéro de facture:</b> {invoice['invoice_number']}<br/>
        <b>Date:</b> {datetime.strptime(invoice['invoice_date'], '%Y-%m-%d').strftime('%d/%m/%Y')}<br/>
        <b>Échéance:</b> {datetime.strptime(invoice['due_date'], '%Y-%m-%d').strftime('%d/%m/%Y')}<br/>
        <b>Patient:</b> {invoice['patient_name']}<br/>
        <b>Statut:</b> {invoice['status'].upper()}
        """
        content.append(Paragraph(invoice_info, normal_style))
        content.append(Spacer(1, 20))
        
        # Treatments table
        content.append(Paragraph("Détail des traitements", heading_style))
        
        # Get invoice items
        invoice_items = practice_db.get_invoice_items(invoice_id)
        
        if invoice_items:
            # Table headers
            table_data = [['Code TARMED', 'Traitement', 'Qté', 'Prix unitaire', 'Total', 'LAMal']]
            
            # Table rows
            for item in invoice_items:
                row = [
                    item['tarmed_code'],
                    item['treatment_name'],
                    str(item['quantity']),
                    f"{item['unit_price_chf']:.2f} CHF",
                    f"{item['total_price_chf']:.2f} CHF",
                    'Oui' if item['lamal_covered'] else 'Non'
                ]
                table_data.append(row)
            
            # Create table
            treatments_table = Table(table_data, colWidths=[1*inch, 2.5*inch, 0.5*inch, 1*inch, 1*inch, 0.8*inch])
            treatments_table.setStyle(TableStyle([
                # Header styling
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2c5aa0')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                
                # Body styling
                ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                ('FONTSIZE', (0, 1), (-1, -1), 9),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.beige, colors.white]),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ('LEFTPADDING', (0, 0), (-1, -1), 6),
                ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                ('TOPPADDING', (0, 0), (-1, -1), 8),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
            ]))
            
            content.append(treatments_table)
            content.append(Spacer(1, 20))
        
        # Totals
        content.append(Paragraph("Récapitulatif", heading_style))
        totals_text = f"""
        <b>Total:</b> {invoice['total_amount_chf']:.2f} CHF<br/>
        <b>Prise en charge LAMal:</b> {invoice['lamal_amount_chf']:.2f} CHF<br/>
        <b>Assurance complémentaire:</b> {invoice['insurance_amount_chf']:.2f} CHF<br/>
        <b>À payer par le patient:</b> {invoice['patient_amount_chf']:.2f} CHF<br/>
        <b>Déjà payé:</b> {invoice['paid_amount_chf']:.2f} CHF<br/>
        <b>Reste à payer:</b> {(invoice['patient_amount_chf'] - invoice['paid_amount_chf']):.2f} CHF
        """
        content.append(Paragraph(totals_text, normal_style))
        content.append(Spacer(1, 30))
        
        # Footer
        footer_text = f"""
        <i>Facture générée le {datetime.now().strftime('%d/%m/%Y à %H:%M')}</i><br/>
        <i>Assistant: Dental AI - Intelligence Artificielle Dentaire</i>
        """
        content.append(Paragraph(footer_text, normal_style))
        
        # Build PDF
        doc.build(content)
        
        # Get PDF data
        pdf_data = buffer.getvalue()
        buffer.close()
        
        # Create response
        response = app.response_class(
            pdf_data,
            mimetype='application/pdf',
            headers={
                'Content-Disposition': f'attachment; filename=facture-{invoice["invoice_number"]}.pdf',
                'Content-Length': len(pdf_data)
            }
        )
        
        return response
        
    except ImportError:
        return jsonify({
            'success': False,
            'error': 'PDF generation library not installed. Please install reportlab.'
        }), 500
    except Exception as e:
        print(f"❌ Error generating invoice PDF: {e}")
        return jsonify({
            'success': False,
            'error': f'Erreur lors de la génération du PDF: {str(e)}'
        }), 500

@app.route('/api/invoices/<invoice_id>', methods=['GET'])
def get_invoice_details(invoice_id):
    """Get detailed information about a specific invoice"""
    try:
        # Get invoice data
        invoices = practice_db.get_invoices(invoice_id=invoice_id)
        if not invoices:
            return jsonify({'success': False, 'error': 'Facture non trouvée'}), 404
        
        invoice = invoices[0]
        
        # Get invoice items
        invoice_items = practice_db.get_invoice_items(invoice_id)
        invoice['items'] = invoice_items
        
        return jsonify({
            'success': True,
            'invoice': invoice
        })
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/api/devis/<devis_id>', methods=['DELETE'])
def delete_devis(devis_id):
    """Delete a devis"""
    try:
        success = practice_db.delete_devis(devis_id)
        
        if success:
            return jsonify({
                'success': True,
                'message': 'Devis supprimé avec succès'
            })
        else:
            return jsonify({
                'success': False,
                'error': 'Devis non trouvé'
            }), 404
            
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/api/invoices/<invoice_id>', methods=['GET', 'DELETE'])
def manage_invoice(invoice_id):
    """Manage specific invoice - GET for details, DELETE to remove"""
    if request.method == 'GET':
        try:
            # Get invoice data
            invoices = practice_db.get_invoices(invoice_id=invoice_id)
            if not invoices:
                return jsonify({'success': False, 'error': 'Facture non trouvée'}), 404
            
            invoice = invoices[0]
            
            # Get invoice items
            invoice_items = practice_db.get_invoice_items(invoice_id)
            invoice['items'] = invoice_items
            
            return jsonify({
                'success': True,
                'invoice': invoice
            })
            
        except Exception as e:
            return jsonify({
                'success': False,
                'error': str(e)
            }), 500
    
    elif request.method == 'DELETE':
        try:
            success = practice_db.delete_invoice(invoice_id)
            
            if success:
                return jsonify({
                    'success': True,
                    'message': 'Facture supprimée avec succès'
                })
            else:
                return jsonify({
                    'success': False,
                    'error': 'Facture non trouvée'
                }), 404
                
        except Exception as e:
            return jsonify({
                'success': False,
                'error': str(e)
            }), 500

@app.route('/api/generate-patient-education', methods=['POST'])
def generate_patient_education():
    """Generate patient education document using specialized LLM"""
    try:
        data = request.get_json()
        patient_id = data.get('patient_id')
        treatment_plan = data.get('treatment_plan')
        
        if not patient_id or not treatment_plan:
            return jsonify({
                'success': False,
                'error': 'Patient ID and treatment plan are required'
            }), 400
        
        # Get patient information
        patient = practice_db.get_patient(patient_id)
        if not patient:
            return jsonify({
                'success': False,
                'error': 'Patient not found'
            }), 404
        
        # Get the patient-education specialized LLM
        education_llm = specialized_llms.get('patient-education')
        if not education_llm:
            return jsonify({
                'success': False,
                'error': 'Patient education LLM not available'
            }), 500
        
        # Build education prompt
        patient_name = f"{patient['first_name']} {patient['last_name']}"
        
        # Calculate age safely, handling empty or missing birth dates
        age_info = ""
        if patient.get('birth_date') and patient['birth_date'].strip():
            try:
                age = datetime.now().year - datetime.strptime(patient['birth_date'], '%Y-%m-%d').year
                age_info = f" (âge: {age} ans)"
            except (ValueError, TypeError):
                age_info = ""
        
        education_prompt = f"""
        Créez un document éducatif pour le patient {patient_name}{age_info}.
        
        DIAGNOSTIC: {treatment_plan.get('consultation_text', 'Non spécifié')}
        
        PLAN DE TRAITEMENT:
        {chr(10).join([f"• RDV {step.get('rdv', i+1)}: {step.get('traitement', 'Non spécifié')} ({step.get('duree', 'Durée non spécifiée')})" 
                       for i, step in enumerate(treatment_plan.get('treatment_sequence', []))])}
        
        Créez un document éducatif complet qui explique:
        1. Le diagnostic en termes simples
        2. Pourquoi ce traitement est nécessaire
        3. Comment se déroule chaque étape
        4. Les bénéfices attendus
        5. Les soins post-traitement
        6. Les conseils de prévention
        
        Utilisez un ton rassurant et professionnel, évitez le jargon médical.
        """
        
        # Generate education content
        response_data = education_llm.generate_response(education_prompt)
        
        if not response_data.get('success', False):
            return jsonify({
                'success': False,
                'error': f'Error generating education content: {response_data.get("error", "Unknown error")}'
            }), 500
        
        education_content = response_data.get('response', '')
        
        # Format the content for HTML display
        formatted_content = education_content.replace('\n', '<br>')
        
        return jsonify({
            'success': True,
            'education_content': formatted_content,
            'patient_name': patient_name
        })
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/api/save-patient-education', methods=['POST'])
def save_patient_education():
    """Save patient education document"""
    try:
        data = request.get_json()
        patient_id = data.get('patient_id')
        education_content = data.get('education_content')
        education_title = data.get('education_title')
        treatment_plan_id = data.get('treatment_plan_id')
        
        if not patient_id or not education_content:
            return jsonify({
                'success': False,
                'error': 'Patient ID and education content are required'
            }), 400
        
        # Extract title from content if not provided
        if not education_title:
            # Try to extract title from the first line or use a default
            lines = education_content.split('\n')
            first_line = lines[0].strip() if lines else ''
            if first_line and len(first_line) < 100:
                education_title = first_line.replace('<h1>', '').replace('</h1>', '').replace('<h2>', '').replace('</h2>', '').strip()
            else:
                education_title = f"Document éducatif - {datetime.now().strftime('%d/%m/%Y')}"
        
        # Save to database
        education_id = practice_db.create_patient_education(
            patient_id=patient_id,
            education_content=education_content,
            education_title=education_title,
            treatment_plan_id=treatment_plan_id
        )
        
        if education_id:
            return jsonify({
                'success': True,
                'education_id': education_id,
                'message': 'Document éducatif sauvegardé avec succès'
            })
        else:
            return jsonify({
                'success': False,
                'error': 'Erreur lors de la sauvegarde du document éducatif'
            }), 500
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

@app.route('/api/download-patient-education', methods=['POST'])
def download_patient_education():
    """Generate and download patient education PDF"""
    try:
        from reportlab.lib.pagesizes import letter, A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib import colors
        from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
        from io import BytesIO
        import html2text
        
        data = request.get_json()
        patient_id = data.get('patient_id')
        patient_name = data.get('patient_name')
        education_content = data.get('education_content')
        
        if not patient_id or not education_content:
            return jsonify({'error': 'Patient ID and education content are required'}), 400
        
        # Create PDF
        buffer = BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4, 
                              rightMargin=72, leftMargin=72, 
                              topMargin=72, bottomMargin=18)
        
        # Get styles
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=30,
            alignment=TA_CENTER,
            textColor=colors.HexColor('#2c5aa0')
        )
        
        subtitle_style = ParagraphStyle(
            'CustomSubtitle',
            parent=styles['Heading2'],
            fontSize=14,
            spaceAfter=12,
            spaceBefore=20,
            textColor=colors.HexColor('#2c5aa0')
        )
        
        body_style = ParagraphStyle(
            'CustomBody',
            parent=styles['Normal'],
            fontSize=11,
            spaceAfter=12,
            alignment=TA_JUSTIFY,
            leftIndent=0,
            rightIndent=0
        )
        
        # Build PDF content
        story = []
        
        # Title
        story.append(Paragraph("Document Éducatif Patient", title_style))
        story.append(Spacer(1, 12))
        
        # Patient info
        story.append(Paragraph(f"<b>Patient:</b> {patient_name}", body_style))
        story.append(Paragraph(f"<b>Date:</b> {datetime.now().strftime('%d/%m/%Y')}", body_style))
        story.append(Spacer(1, 20))
        
        # Convert HTML content to plain text and then back to paragraphs
        h = html2text.HTML2Text()
        h.ignore_links = True
        h.ignore_images = True
        plain_text = h.handle(education_content)
        
        # Split content into paragraphs and format
        paragraphs = plain_text.split('\n\n')
        
        for para in paragraphs:
            if para.strip():
                # Check if it's a title (starts with #)
                if para.strip().startswith('#'):
                    title_text = para.strip().replace('#', '').strip()
                    story.append(Paragraph(title_text, subtitle_style))
                elif para.strip().startswith('•') or para.strip().startswith('-'):
                    # Handle bullet points
                    story.append(Paragraph(para.strip(), body_style))
                else:
                    story.append(Paragraph(para.strip(), body_style))
        
        # Footer
        story.append(Spacer(1, 30))
        footer_style = ParagraphStyle(
            'Footer',
            parent=styles['Normal'],
            fontSize=9,
            alignment=TA_CENTER,
            textColor=colors.grey
        )
        story.append(Paragraph("Ce document a été généré par votre cabinet dentaire", footer_style))
        story.append(Paragraph("Pour toute question, n'hésitez pas à nous contacter", footer_style))
        
        # Build PDF
        doc.build(story)
        
        # Return PDF
        buffer.seek(0)
        return send_file(
            buffer,
            as_attachment=True,
            download_name=f"Education_{patient_name.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.pdf",
            mimetype='application/pdf'
        )
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/patient-education', methods=['GET'])
def get_patient_education():
    """Get patient education documents"""
    try:
        patient_id = request.args.get('patient_id')
        education_id = request.args.get('education_id')
        
        if education_id:
            # Get specific education document
            education_doc = practice_db.get_patient_education(education_id=education_id)
            if education_doc:
                return jsonify({
                    'success': True,
                    'education_document': education_doc
                })
            else:
                return jsonify({
                    'success': False,
                    'error': 'Document éducatif non trouvé'
                }), 404
        
        elif patient_id:
            # Get all education documents for a patient
            education_docs = practice_db.get_patient_education(patient_id=patient_id)
            return jsonify({
                'success': True,
                'education_documents': education_docs,
                'count': len(education_docs)
            })
        
        else:
            # Get all education documents
            education_docs = practice_db.get_patient_education()
            return jsonify({
                'success': True,
                'education_documents': education_docs,
                'count': len(education_docs)
            })
        
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        }), 500

# PowerPoint Generation System
# Enhanced treatment mappings - French dental terms to actions
TREATMENT_MAPPINGS = {
    # Color changes (crown, onlay, veneer)
    'color_treatments': {
        'cc': 'Couronne céramique',
        'cci': 'Couronne sur implant',
        'couronne': 'Couronne céramique',
        'couronne ceramique': 'Couronne céramique',
        'couronne sur implant': 'Couronne sur implant',
        'onlay': 'Onlay',
        'facette': 'Facette céramique',
        'facette ceramique': 'Facette céramique',
        'veneer': 'Facette céramique',
        'f': 'Facette céramique',
        'cpr': 'Couronne céramique',
        'o': 'Onlay',
        'ceram': 'Facette céramique',
        'ceramique': 'Facette céramique'
    },
    # Icon treatments
    'icon_treatments': {
        'dem': 'Dévitalisation',
        'dém': 'Dévitalisation',
        'devitalisation': 'Dévitalisation',
        'dévitalisation': 'Dévitalisation',
        'tr': 'Traitement endodontique',
        'traitement radiculaire': 'Traitement endodontique',
        'traitement endodontique': 'Traitement endodontique',
        'endo': 'Traitement endodontique',
        'endodontie': 'Traitement endodontique',
        'tenons': 'Tenons',
        'tenon': 'Tenons',
        'ma': 'Moignon adhésif',
        'moignon adhesif': 'Moignon adhésif',
        'moignon adhésif': 'Moignon adhésif',
        'extraction': 'Extraction',
        'ext': 'Extraction',
        'av': 'Extraction',
        'avulsion': 'Extraction',
        'abl': 'Extraction',
        'ablation': 'Extraction',
        'implant': 'Pose d\'implant',
        'imp': 'Pose d\'implant',
        'pose i': 'Pose d\'implant',
        'pose d\'implant': 'Pose d\'implant',
        'pose implant': 'Pose d\'implant',
        'curtage': 'Curetage',
        'curetage': 'Curetage',
        'seance': 'Séance',
        'séance': 'Séance',
        'bnv': 'Blanchissement interne',
        'blanchiment interne': 'Blanchissement interne',
        'blanchissement interne': 'Blanchissement interne',
        'blanch': 'Blanchissement interne',
        'blanchiment': 'Blanchissement interne',
        'gbr': 'Greffe osseuse',
        'greffe osseuse': 'Greffe osseuse',
        'gc': 'Greffe gingivale',
        'greffe gingivale': 'Greffe gingivale',
        'sl': 'Sinus lift',
        'sinus lift': 'Sinus lift',
        'det': 'Détartrage',
        'hd': 'Détartrage',
        'détartrage': 'Détartrage',
        'detartrage': 'Détartrage',
        'sf': 'Scellement de fissure',
        'scellement de fissure': 'Scellement de fissure',
        'dds': 'Dent de sagesse',
        'dent de sagesse': 'Dent de sagesse',
        'fil de cont': 'Fil de contention',
        'fil de contention': 'Fil de contention',
        'dem cc': 'Démonter couronne',
        'dém cc': 'Démonter couronne',
        'te': 'Taille empreinte',
        'taille empreinte': 'Taille empreinte',
        'sc': 'Scellement',
        'scellement': 'Scellement',
        'empr': 'Empreinte',
        'empreinte': 'Empreinte',
        'post-op': 'Post opératoire',
        'post op': 'Post opératoire',
        'prov': 'Provisoire',
        'provisoire': 'Provisoire',
        'm': 'Composite mésial',
        'mesial': 'Composite mésial',
        'mésial': 'Composite mésial',
        'd': 'Composite distal',
        'distal': 'Composite distal',
        'mo': 'Composite mésio-occlusal',
        'do': 'Composite occluso-distal',
        'mod': 'Composite mésio-occluso-distal',
        'l': 'Composite lingual',
        'lingual': 'Composite lingual',
        'p': 'Composite palatin',
        'palatin': 'Composite palatin',
        'v': 'Composite vestibulaire',
        'vestibulaire': 'Composite vestibulaire'
    }
}

# Treatment colors
TREATMENT_COLORS = {
    'Couronne céramique': RGBColor(255, 215, 0),  # Gold
    'Couronne sur implant': RGBColor(255, 165, 0),  # Orange
    'Onlay': RGBColor(192, 192, 192),  # Silver
    'Facette céramique': RGBColor(0, 123, 255)  # Blue
}

def parse_tooth_range(tooth_str):
    """Parse tooth ranges like '12-22' or '11 à 22'"""
    tooth_str = tooth_str.strip()
    
    # Handle ranges with dash: 12-22
    if '-' in tooth_str:
        start, end = tooth_str.split('-')
        start, end = int(start.strip()), int(end.strip())
        return list(range(start, end + 1))
    
    # Handle ranges with 'à': 11 à 22
    elif ' à ' in tooth_str:
        start, end = tooth_str.split(' à ')
        start, end = int(start.strip()), int(end.strip())
        return list(range(start, end + 1))
    
    # Single tooth
    else:
        return [int(tooth_str)]

def enhanced_parse_treatment_text(text):
    """Enhanced parsing with better regex"""
    results = []
    
    # Clean the text
    text = text.lower().strip()
    
    # Enhanced regex patterns
    patterns = [
        # Complex pattern: Plan de TT 11 AV + implant + CC; 22 Implant + CC
        r'plan\s+de\s+t+\s+([^;]+)',
        # Simple pattern: 26 dém. CC + dém. tenons + TR
        r'(\d+(?:\s*[-à]\s*\d+)?)\s*[:\s]*([^;]+)',
        # Alternative pattern: Pour la 26: treatments
        r'pour\s+la\s+(\d+(?:\s*[-à]\s*\d+)?)\s*[:\s]*([^;]+)',
    ]
    
    for pattern in patterns:
        matches = re.findall(pattern, text)
        for match in matches:
            if len(match) == 2:
                tooth_part, treatment_part = match
                
                # Parse tooth numbers (handle ranges)
                try:
                    tooth_numbers = parse_tooth_range(tooth_part)
                except ValueError:
                    continue
                
                # Parse treatments
                treatments = [t.strip() for t in re.split(r'[+&,]', treatment_part) if t.strip()]
                
                for tooth_num in tooth_numbers:
                    for treatment in treatments:
                        treatment = treatment.strip()
                        if treatment:
                            # Determine treatment type and normalize name
                            normalized_treatment = None
                            treatment_type = 'icon'  # default
                            
                            # Check color treatments first
                            for key, value in TREATMENT_MAPPINGS['color_treatments'].items():
                                if key in treatment:
                                    normalized_treatment = value
                                    treatment_type = 'color'
                                    break
                            
                            # Check icon treatments if not found in color treatments
                            if not normalized_treatment:
                                for key, value in TREATMENT_MAPPINGS['icon_treatments'].items():
                                    if key in treatment:
                                        normalized_treatment = value
                                        treatment_type = 'icon'
                                        break
                            
                            # If no mapping found, use original treatment
                            if not normalized_treatment:
                                normalized_treatment = treatment.title()
                            
                            results.append({
                                'tooth': str(tooth_num),
                                'treatment': normalized_treatment,
                                'type': treatment_type,
                                'original': treatment
                            })
    
    return results

def is_valid_tooth_number(tooth_number):
    """Validate tooth number according to FDI system"""
    try:
        tooth_num = int(tooth_number)
        # Valid FDI tooth numbers: 11-18, 21-28, 31-38, 41-48
        valid_ranges = [
            (11, 18), (21, 28), (31, 38), (41, 48)
        ]
        
        for start, end in valid_ranges:
            if start <= tooth_num <= end:
                return True
        return False
    except ValueError:
        return False

def find_tooth_element(slide, tooth_number):
    """Find the tooth element in the slide by searching for tooth name patterns"""
    target_names = [
        f"tooth_{tooth_number}",
        f"Tooth_{tooth_number}",
        f"background_tooth_{tooth_number}",
        f"Background_tooth_{tooth_number}",
        f"tooth{tooth_number}",
        f"Tooth{tooth_number}",
        f"dent_{tooth_number}",
        f"Dent_{tooth_number}"
    ]
    
    def search_in_shapes(shapes):
        # First, check direct shapes by name (faster and more accurate)
        for shape in shapes:
            if hasattr(shape, 'name') and shape.name in target_names:
                return shape
        
        # Then check grouped shapes
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                result = search_in_shapes(shape.shapes)
                if result:
                    return result
        
        # Fallback: check if tooth number appears in shape text
        for shape in shapes:
            if hasattr(shape, 'text') and str(tooth_number) in shape.text:
                return shape
        
        return None
    
    return search_in_shapes(slide.shapes)

def apply_color_treatment(slide, tooth_number, treatment):
    """Apply color treatment to a tooth with enhanced error handling"""
    tooth_element = find_tooth_element(slide, tooth_number)
    
    if not tooth_element:
        return False
    
    try:
        # Treatment color mapping
        treatment_colors = {
            'Extraction': RGBColor(255, 0, 0),      # Red
            'Couronne': RGBColor(255, 215, 0),      # Gold
            'Couronne céramique': RGBColor(255, 215, 0),      # Gold
            'Couronne sur implant': RGBColor(255, 165, 0),    # Orange
            'Implant': RGBColor(0, 128, 255),       # Blue
            'Obturation': RGBColor(169, 169, 169),  # Gray
            'Endodontie': RGBColor(255, 165, 0),    # Orange
            'Prothèse': RGBColor(128, 0, 128),      # Purple
            'Bridge': RGBColor(255, 192, 203),      # Pink
            'Facette': RGBColor(255, 255, 255),     # White
            'Facette céramique': RGBColor(0, 123, 255),  # Blue
            'Onlay': RGBColor(192, 192, 192),       # Silver
            'Détartrage': RGBColor(144, 238, 144),  # Light Green
            'Blanchiment': RGBColor(255, 255, 224), # Light Yellow
            'Gingivectomie': RGBColor(255, 20, 147), # Deep Pink
            'Greffe osseuse': RGBColor(222, 184, 135), # Burlywood
            'Sinus lift': RGBColor(205, 133, 63),   # Peru
            'Chirurgie gingivale': RGBColor(255, 105, 180), # Hot Pink
            'Restauration composite': RGBColor(240, 248, 255), # Alice Blue
            'Scellement': RGBColor(192, 192, 192),  # Silver
            'Provisoire': RGBColor(255, 228, 196),  # Bisque
        }
        
        color = treatment_colors.get(treatment, RGBColor(128, 128, 128))  # Default gray
        
        # Try multiple approaches to set the color
        # Approach 1: Direct fill
        if hasattr(tooth_element, 'fill'):
            try:
                tooth_element.fill.solid()
                tooth_element.fill.fore_color.rgb = color
                return True
            except Exception:
                pass
        
        # Approach 2: Try line color if fill doesn't work
        if hasattr(tooth_element, 'line'):
            try:
                tooth_element.line.color.rgb = color
                tooth_element.line.width = Pt(3)
                return True
            except Exception:
                pass
        
        # Approach 3: Try text color if it's a text shape
        if hasattr(tooth_element, 'text_frame'):
            try:
                if tooth_element.text_frame.paragraphs:
                    tooth_element.text_frame.paragraphs[0].font.color.rgb = color
                    return True
            except Exception:
                pass
        
        # Approach 4: For grouped shapes, try to color all sub-shapes
        if hasattr(tooth_element, 'shapes'):
            try:
                colored_count = 0
                for sub_shape in tooth_element.shapes:
                    try:
                        if hasattr(sub_shape, 'fill'):
                            sub_shape.fill.solid()
                            sub_shape.fill.fore_color.rgb = color
                            colored_count += 1
                    except:
                        pass
                if colored_count > 0:
                    return True
            except Exception:
                pass
        
        return False
        
    except Exception as e:
        print(f"Error applying color to tooth {tooth_number}: {e}")
        return False

def apply_multiple_icon_treatments(slide, tooth_number, treatments):
    """Apply multiple icon treatments to a tooth with smart positioning"""
    try:
        tooth_element = find_tooth_element(slide, tooth_number)
        if not tooth_element:
            return [False] * len(treatments)
        
        # Get tooth element position and dimensions
        tooth_left = tooth_element.left
        tooth_top = tooth_element.top
        tooth_width = tooth_element.width
        tooth_height = tooth_element.height
        
        # Calculate positions for multiple icons
        icon_size = Inches(0.25)  # Slightly larger for better visibility
        positions = []
        
        if len(treatments) == 1:
            # Single icon - center
            center_x = tooth_left + tooth_width / 2 - icon_size / 2
            center_y = tooth_top + tooth_height / 2 - icon_size / 2
            positions.append((center_x, center_y))
        elif len(treatments) == 2:
            # Two icons - side by side
            spacing = icon_size * 0.1
            total_width = icon_size * 2 + spacing
            start_x = tooth_left + (tooth_width - total_width) / 2
            center_y = tooth_top + tooth_height / 2 - icon_size / 2
            positions.append((start_x, center_y))
            positions.append((start_x + icon_size + spacing, center_y))
        else:
            # Multiple icons - grid layout
            icons_per_row = 2
            rows = (len(treatments) + icons_per_row - 1) // icons_per_row
            
            row_height = icon_size * 0.8
            total_height = row_height * rows
            start_y = tooth_top + (tooth_height - total_height) / 2
            
            for i, treatment in enumerate(treatments):
                row = i // icons_per_row
                col = i % icons_per_row
                
                if row == rows - 1 and len(treatments) % icons_per_row != 0:
                    # Last row with fewer icons - center them
                    remaining_icons = len(treatments) % icons_per_row
                    col_spacing = tooth_width / (remaining_icons + 1)
                    x = tooth_left + col_spacing * (col + 1) - icon_size / 2
                else:
                    col_spacing = tooth_width / (icons_per_row + 1)
                    x = tooth_left + col_spacing * (col + 1) - icon_size / 2
                
                y = start_y + row * row_height
                positions.append((x, y))
        
        # Apply each treatment
        results = []
        for i, treatment in enumerate(treatments):
            try:
                x, y = positions[i]
                success = False
                
                # Try to find and use an icon file first
                icon_path = get_icon_path(treatment)
                
                if icon_path and os.path.exists(icon_path):
                    try:
                        print(f"Adding icon for {treatment}: {icon_path}")
                        # Add the icon image
                        pic = slide.shapes.add_picture(
                            icon_path,
                            x,
                            y,
                            width=icon_size,
                            height=icon_size
                        )
                        success = True
                        print(f"✅ Successfully added icon for {treatment}")
                    except Exception as e:
                        print(f"❌ Error adding icon image for {treatment}: {e}")
                        # Fall back to text if image fails
                        pass
                
                if not success:
                    print(f"Using text fallback for {treatment}")
                    # Fallback to text
                    text_box = slide.shapes.add_textbox(x, y, icon_size, icon_size)
                    
                    # Treatment symbols mapping
                    symbols = {
                        'Extraction': 'EX',
                        'Couronne': 'C',
                        'Implant': 'I',
                        'Obturation': 'O',
                        'Endodontie': 'E',
                        'Traitement endodontique': 'TR',
                        'Dévitalisation': 'DÉM',
                        'Prothèse': 'P',
                        'Bridge': 'B',
                        'Facette': 'F',
                        'Détartrage': 'DET',
                        'Blanchiment': 'BL',
                        'Blanchissement interne': 'BNV',
                        'Gingivectomie': 'G',
                        'Greffe osseuse': 'GBR',
                        'Sinus lift': 'SL',
                        'Chirurgie gingivale': 'CG',
                        'Greffe gingivale': 'GC',
                        'Restauration composite': 'RC',
                        'Moignon adhésif': 'MA',
                        'Tenons': 'T',
                        'Curetage': 'C',
                        'Séance': 'S',
                        'Démonter couronne': 'DC',
                        'Démonter tenon': 'DT',
                        'Taille empreinte': 'TE',
                        'Scellement': 'SC',
                        'Empreinte': 'E',
                        'Post opératoire': 'PO',
                        'Dent de sagesse': 'DS',
                        'Fil de contention': 'FC',
                        'Scellement de fissure': 'SF',
                        'Dimension verticale occlusion': 'DVO',
                        'Augmenter DVO': 'A+',
                        'Provisoire': 'P',
                        'Traitement': 'T',
                        'Composite mésial': 'M',
                        'Composite distal': 'D',
                        'Composite mésio-occlusal': 'MO',
                        'Composite occluso-distal': 'OD',
                        'Composite mésio-occluso-distal': 'MOD',
                        'Composite interproximal': 'IP',
                        'Composite lingual': 'L',
                        'Composite palatin': 'P',
                        'Composite vestibulaire': 'V'
                    }
                    
                    text_frame = text_box.text_frame
                    text_frame.text = symbols.get(treatment, treatment[:3].upper())
                    text_frame.margin_left = 0
                    text_frame.margin_right = 0
                    text_frame.margin_top = 0
                    text_frame.margin_bottom = 0
                    
                    # Style the text (smaller for multiple icons)
                    paragraph = text_frame.paragraphs[0]
                    font_size = Pt(10) if len(treatments) == 1 else Pt(8)
                    paragraph.font.size = font_size
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(255, 0, 0)  # Red
                    
                    # Center the text
                    from pptx.enum.text import PP_ALIGN
                    paragraph.alignment = PP_ALIGN.CENTER
                    
                    success = True
                
                results.append(success)
                
            except Exception as e:
                print(f"Icon treatment error for tooth {tooth_number}, treatment {treatment}: {e}")
                results.append(False)
        
        return results
        
    except Exception as e:
        print(f"Multiple icon treatment error for tooth {tooth_number}: {e}")
        return [False] * len(treatments)

def apply_icon_treatment(slide, tooth_number, treatment):
    """Apply single icon treatment (wrapper for backward compatibility)"""
    results = apply_multiple_icon_treatments(slide, tooth_number, [treatment])
    return results[0] if results else False

def debug_slide_shapes(slide):
    """Debug function to list all shapes in a slide"""
    shapes_info = []
    
    def collect_shapes(shapes, level=0):
        indent = "  " * level
        for i, shape in enumerate(shapes):
            shape_info = {
                'index': i,
                'name': getattr(shape, 'name', 'No name'),
                'type': shape.shape_type,
                'level': level
            }
            
            # Check if it has text
            if hasattr(shape, 'text'):
                shape_info['text'] = shape.text[:50] if shape.text else ''
            
            shapes_info.append(shape_info)
            print(f"{indent}Shape {i}: {shape_info['name']} (type: {shape_info['type']})")
            
            # If it's a group, recurse
            if hasattr(shape, 'shapes'):
                collect_shapes(shape.shapes, level + 1)
    
    collect_shapes(slide.shapes)
    return shapes_info

def process_powerpoint_treatments(treatments):
    """Process the PowerPoint with the given treatments"""
    try:
        # Load the PowerPoint template
        template_path = 'plan.pptx'
        if not os.path.exists(template_path):
            return None, "Template PowerPoint file not found"
        
        prs = Presentation(template_path)
        
        results = []
        
        # Apply treatments to the first slide (assuming dental chart is on first slide)
        if prs.slides:
            slide = prs.slides[0]
            
            # Debug: Print all shapes in the slide
            print("=== DEBUG: Slide shapes ===")
            debug_slide_shapes(slide)
            print("=== END DEBUG ===")
            
            # Group treatments by tooth and type
            tooth_treatments = {}
            
            for treatment in treatments:
                tooth = treatment['tooth']
                
                # First check if it's a valid tooth number
                if not is_valid_tooth_number(tooth):
                    results.append({
                        'tooth': tooth,
                        'treatment': treatment['treatment'],
                        'success': False,
                        'error': f"Numéro de dent invalide ({tooth} n'existe pas dans le système FDI)"
                    })
                    continue
                
                if tooth not in tooth_treatments:
                    tooth_treatments[tooth] = {'color': [], 'icon': []}
                
                tooth_treatments[tooth][treatment['type']].append(treatment)
            
            # Process each tooth's treatments
            for tooth, treatments_by_type in tooth_treatments.items():
                print(f"Processing tooth {tooth}...")
                
                # Try to find the tooth element for debugging
                tooth_element = find_tooth_element(slide, tooth)
                if tooth_element:
                    print(f"Found tooth element for {tooth}: {tooth_element.name}")
                else:
                    print(f"Could not find tooth element for {tooth}")
                
                # Apply color treatments first (they don't stack)
                for color_treatment in treatments_by_type['color']:
                    success = apply_color_treatment(slide, tooth, color_treatment['treatment'])
                    results.append({
                        'tooth': tooth,
                        'treatment': color_treatment['treatment'],
                        'success': success,
                        'error': f"Élément tooth_{tooth} non trouvé dans le PowerPoint" if not success else None
                    })
                
                # Apply all icon treatments for this tooth at once (with smart positioning)
                if treatments_by_type['icon']:
                    icon_treatments = [t['treatment'] for t in treatments_by_type['icon']]
                    success_list = apply_multiple_icon_treatments(slide, tooth, icon_treatments)
                    
                    for i, (icon_treatment, success) in enumerate(zip(treatments_by_type['icon'], success_list)):
                        results.append({
                            'tooth': tooth,
                            'treatment': icon_treatment['treatment'],
                            'success': success,
                            'error': f"Impossible d'ajouter l'icône sur la dent {tooth}" if not success else None
                        })
        
        # Save the modified presentation
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_filename = f"plan_modified_{timestamp}.pptx"
        
        # Create temp directory if it doesn't exist
        os.makedirs('temp', exist_ok=True)
        output_path = os.path.join('temp', output_filename)
        
        prs.save(output_path)
        
        return output_filename, results
        
    except Exception as e:
        print(f"PowerPoint processing error: {e}")
        return None, str(e)

@app.route('/api/process-powerpoint', methods=['POST'])
def process_powerpoint():
    """Process PowerPoint generation request"""
    try:
        data = request.get_json()
        text = data.get('text', '').strip()
        
        if not text:
            return jsonify({'success': False, 'error': 'Aucun texte fourni'}), 400
        
        # Parse the treatment text
        treatments = enhanced_parse_treatment_text(text)
        
        if not treatments:
            return jsonify({'success': False, 'error': 'Aucun traitement reconnu dans le texte'}), 400
        
        # Process the PowerPoint
        output_file, results = process_powerpoint_treatments(treatments)
        
        if output_file:
            return jsonify({
                'success': True,
                'treatments': results,
                'output_file': output_file
            })
        else:
            return jsonify({'success': False, 'error': results}), 500
            
    except Exception as e:
        print(f"Error in process_powerpoint: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/download-powerpoint/<filename>')
def download_powerpoint(filename):
    """Download the generated PowerPoint file"""
    try:
        file_path = os.path.join('temp', filename)
        if not os.path.exists(file_path):
            return jsonify({'error': 'File not found'}), 404
        
        return send_file(
            file_path,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception as e:
        print(f"Error downloading PowerPoint: {e}")
        return jsonify({'error': str(e)}), 500

def get_icon_path(treatment):
    """Get the path to the icon file for a treatment"""
    # Treatment to icon filename mapping
    icon_mappings = {
        # Main treatments
        'Blanchissement interne': 'bnv.png',
        'Traitement endodontique': 'tr.png',
        'Dévitalisation': 'tr.png',  # Use TR icon for dévitalisation
        'Extraction': 'extraction.png',
        'Pose d\'implant': 'implant.png',
        'Implant': 'implant.png',
        'Greffe osseuse': 'gbr.png',
        'Greffe gingivale': 'gc.png',
        'Moignon adhésif': 'ma.png',
        'Détartrage': 'det.png',
        
        # Additional mappings for common treatments
        'Endodontie': 'tr.png',
        'Avulsion': 'extraction.png',
        'Curetage': 'gc.png',  # Use GC icon for curetage
        'Séance': 'det.png',   # Use DET icon for general séance
        
        # Composite treatments (could use a general composite icon if available)
        'Restauration composite': 'ma.png',  # Use MA as fallback
        'Composite mésial': 'ma.png',
        'Composite distal': 'ma.png',
        'Composite mésio-occlusal': 'ma.png',
        'Composite occluso-distal': 'ma.png',
        'Composite mésio-occluso-distal': 'ma.png',
        'Composite interproximal': 'ma.png',
        'Composite lingual': 'ma.png',
        'Composite palatin': 'ma.png',
        'Composite vestibulaire': 'ma.png',
        
        # Other treatments that might map to existing icons
        'Tenons': 'ma.png',
        'Chirurgie gingivale': 'gc.png',
        'Sinus lift': 'gbr.png',  # Use GBR for sinus lift
        'Greffe d\'os': 'gbr.png',
        'Greffe de gencive': 'gc.png',
    }
    
    # First try direct mapping
    if treatment in icon_mappings:
        icon_path = f"static/icons/{icon_mappings[treatment]}"
        if os.path.exists(icon_path):
            return icon_path
    
    # Try to create a safe filename from treatment name
    safe_name = re.sub(r'[^\w\s-]', '', treatment.lower())
    safe_name = re.sub(r'[-\s]+', '_', safe_name)
    
    # Common icon file extensions
    extensions = ['.png', '.jpg', '.jpeg', '.gif', '.bmp']
    
    for ext in extensions:
        icon_path = f"static/icons/{safe_name}{ext}"
        if os.path.exists(icon_path):
            return icon_path
    
    # Try common abbreviations
    abbreviations = {
        'Traitement endodontique': 'tr',
        'Dévitalisation': 'tr',
        'Extraction': 'extraction',
        'Pose d\'implant': 'implant',
        'Implant': 'implant',
        'Greffe osseuse': 'gbr',
        'Greffe gingivale': 'gc',
        'Blanchissement interne': 'bnv',
        'Chirurgie gingivale': 'gc',
        'Restauration composite': 'ma',
        'Moignon adhésif': 'ma',
        'Tenons': 'ma',
        'Curetage': 'gc',
        'Détartrage': 'det',
        'Séance': 'det',
    }
    
    if treatment in abbreviations:
        abbrev = abbreviations[treatment]
        for ext in extensions:
            icon_path = f"static/icons/{abbrev}{ext}"
            if os.path.exists(icon_path):
                return icon_path
    
    return None

class IntelligentScheduler:
    """AI-powered intelligent scheduling system for dental treatments"""
    
    def __init__(self, client: OpenAI, practice_db: PracticeDatabase):
        self.client = client
        self.practice_db = practice_db
        self.dentist_preferences = self.get_dentist_preferences()
        self.treatment_rules = self.get_treatment_scheduling_rules()
    
    def get_dentist_preferences(self):
        """Get dentist scheduling preferences - can be made configurable later"""
        return {
            "working_days": ["monday", "tuesday", "wednesday", "thursday", "friday"],
            "working_hours": {
                "start": "08:00",
                "end": "18:00",
                "lunch_break": {"start": "12:00", "end": "13:00"}
            },
            "preferred_schedules": {
                "surgeries": "morning",  # Prefer surgeries in the morning
                "consultations": "afternoon",  # Prefer consultations in the afternoon
                "long_procedures": "morning",  # Long procedures in the morning when fresh
                "follow_ups": "afternoon"  # Follow-ups in the afternoon
            },
            "time_preferences": {
                "first_appointment": "09:00",  # Prefer first appointments at 9 AM
                "last_appointment": "17:00",   # Last appointment at 5 PM
                "buffer_time": 15,  # 15 minutes buffer between appointments
                "emergency_slots": ["08:00", "12:00", "17:00"]  # Reserved emergency slots
            },
            "treatment_spacing": {
                "surgical": {"min_days": 7, "max_days": 14},  # Surgery follow-ups
                "endodontic": {"min_days": 3, "max_days": 10},  # Root canal follow-ups
                "prosthetic": {"min_days": 14, "max_days": 21},  # Prosthetic work
                "routine": {"min_days": 7, "max_days": 30}  # Routine treatments
            },
            "patient_considerations": {
                "elderly": "morning",  # Elderly patients prefer morning
                "children": "afternoon",  # Children after school
                "working_adults": "early_morning_or_evening"
            }
        }
    
    def get_treatment_scheduling_rules(self):
        """Define intelligent rules for different treatment types"""
        return {
            "surgical_treatments": {
                "keywords": ["extraction", "chirurgie", "implant", "greffe", "avulsion", "sinus lift"],
                "preferred_time": "morning",
                "duration_buffer": 30,  # Extra 30 minutes for surgical procedures
                "post_op_required": True,
                "avoid_friday": True  # Avoid Friday surgeries for weekend recovery
            },
            "endodontic_treatments": {
                "keywords": ["endodontie", "dévitalisation", "traitement endodontique", "pulpotomie"],
                "preferred_time": "morning",
                "duration_buffer": 15,
                "multiple_sessions": True,
                "session_spacing": "3-7 days"
            },
            "prosthetic_treatments": {
                "keywords": ["couronne", "bridge", "prothèse", "facette", "onlay", "inlay"],
                "preferred_time": "morning",
                "duration_buffer": 15,
                "multiple_sessions": True,
                "session_spacing": "14-21 days"
            },
            "routine_treatments": {
                "keywords": ["détartrage", "obturation", "composite", "polissage", "contrôle"],
                "preferred_time": "afternoon",
                "duration_buffer": 10,
                "flexible_scheduling": True
            },
            "emergency_treatments": {
                "keywords": ["urgence", "douleur", "abcès", "trauma", "fracture"],
                "preferred_time": "any",
                "priority": "high",
                "duration_buffer": 20
            }
        }
    
    def classify_treatment(self, treatment_name: str) -> dict:
        """Classify a treatment and return its scheduling properties"""
        treatment_lower = treatment_name.lower()
        
        for category, rules in self.treatment_rules.items():
            if any(keyword in treatment_lower for keyword in rules["keywords"]):
                return {
                    "category": category,
                    "preferred_time": rules["preferred_time"],
                    "duration_buffer": rules.get("duration_buffer", 10),
                    "special_requirements": rules
                }
        
        # Default classification
        return {
            "category": "routine_treatments",
            "preferred_time": "afternoon",
            "duration_buffer": 10,
            "special_requirements": self.treatment_rules["routine_treatments"]
        }
    
    def get_patient_preferences(self, patient_id: str) -> dict:
        """Get patient-specific scheduling preferences"""
        patient = self.practice_db.get_patient(patient_id)
        if not patient:
            return {}
        
        # Calculate age if birth_date is available
        age = None
        if patient.get('birth_date'):
            try:
                birth_date = datetime.strptime(patient['birth_date'], '%Y-%m-%d')
                age = (datetime.now() - birth_date).days // 365
            except:
                pass
        
        preferences = {}
        
        # Age-based preferences
        if age:
            if age >= 65:
                preferences["preferred_time"] = "morning"
                preferences["reason"] = "Elderly patient - morning preferred"
            elif age <= 18:
                preferences["preferred_time"] = "afternoon"
                preferences["reason"] = "Young patient - afternoon preferred"
            elif 25 <= age <= 55:
                preferences["preferred_time"] = "early_morning_or_evening"
                preferences["reason"] = "Working adult - early morning or evening preferred"
        
        return preferences
    
    def analyze_current_schedule(self, date: str) -> dict:
        """Analyze the current schedule for a specific date"""
        appointments = self.practice_db.get_appointments_for_date(date)
        
        analysis = {
            "total_appointments": len(appointments),
            "morning_load": 0,
            "afternoon_load": 0,
            "surgical_count": 0,
            "routine_count": 0,
            "available_slots": [],
            "recommendations": []
        }
        
        for apt in appointments:
            hour = int(apt['appointment_time'].split(':')[0])
            if hour < 12:
                analysis["morning_load"] += 1
            else:
                analysis["afternoon_load"] += 1
            
            # Classify existing appointments
            treatment_class = self.classify_treatment(apt.get('treatment_type', ''))
            if treatment_class["category"] in ["surgical_treatments", "endodontic_treatments"]:
                analysis["surgical_count"] += 1
            else:
                analysis["routine_count"] += 1
        
        # Get available slots
        analysis["available_slots"] = self.practice_db.get_available_slots(date, 60)
        
        # Generate recommendations
        if analysis["morning_load"] > analysis["afternoon_load"] + 2:
            analysis["recommendations"].append("Morning is heavily loaded - consider afternoon scheduling")
        if analysis["surgical_count"] >= 3:
            analysis["recommendations"].append("High surgical load - avoid additional surgeries")
        
        return analysis
    
    def generate_intelligent_schedule(self, patient_id: str, treatment_sequence: list, start_date: str) -> dict:
        """Generate an intelligent schedule using LLM analysis"""
        
        try:
            print(f"🔍 DEBUG: Starting intelligent schedule generation")
            print(f"🔍 DEBUG: patient_id = {patient_id} (type: {type(patient_id)})")
            print(f"🔍 DEBUG: treatment_sequence = {treatment_sequence} (type: {type(treatment_sequence)})")
            print(f"🔍 DEBUG: start_date = {start_date} (type: {type(start_date)})")
            
            # Get patient preferences
            patient_prefs = self.get_patient_preferences(patient_id)
            print(f"🔍 DEBUG: patient_prefs = {patient_prefs} (type: {type(patient_prefs)})")
            
            patient = self.practice_db.get_patient(patient_id)
            print(f"🔍 DEBUG: patient = {patient} (type: {type(patient)})")
            
            # Analyze treatments
            treatment_analysis = []
            for i, treatment in enumerate(treatment_sequence):
                print(f"🔍 DEBUG: Processing treatment {i}: {treatment} (type: {type(treatment)})")
                
                classification = self.classify_treatment(treatment.get('traitement', ''))
                print(f"🔍 DEBUG: classification = {classification} (type: {type(classification)})")
                
                treatment_analysis.append({
                    "step": i + 1,
                    "treatment": treatment.get('traitement', ''),
                    "duration": treatment.get('duree', '60 min'),
                    "classification": classification,
                    "original_delay": treatment.get('delai', '1 semaine')
                })
            
            # Prepare context for LLM
            context = {
                "patient_info": {
                    "name": f"{patient.get('first_name', '') if patient else ''} {patient.get('last_name', '') if patient else ''}",
                    "age": patient_prefs.get("age"),
                    "preferences": patient_prefs
                },
                "dentist_preferences": self.dentist_preferences,
                "treatment_analysis": treatment_analysis,
                "start_date": start_date
            }
            
            print(f"🔍 DEBUG: LLM context prepared")
            
            # Generate intelligent scheduling with LLM
            llm_response = self.get_llm_scheduling_recommendations(context)
            print(f"🔍 DEBUG: LLM response = {llm_response} (type: {type(llm_response)})")
            
            # Apply LLM recommendations to create optimized schedule
            optimized_schedule = self.apply_llm_recommendations(
                treatment_sequence, 
                llm_response, 
                start_date, 
                patient_id
            )
            
            print(f"🔍 DEBUG: Optimized schedule generated successfully")
            return optimized_schedule
            
        except Exception as e:
            print(f"❌ ERROR in generate_intelligent_schedule: {str(e)}")
            print(f"❌ ERROR type: {type(e)}")
            import traceback
            traceback.print_exc()
            raise e
    
    def get_llm_scheduling_recommendations(self, context: dict) -> dict:
        """Use LLM to generate intelligent scheduling recommendations"""
        
        system_prompt = """Tu es un assistant intelligent de planification dentaire. 
        Ton rôle est d'analyser les traitements dentaires et de proposer un planning optimal 
        en tenant compte des préférences du dentiste, du patient, et des bonnes pratiques cliniques.

        RÈGLES DE PLANIFICATION:
        1. Chirurgies (extractions, implants, greffes) → Matinée, éviter vendredi
        2. Endodontie → Matinée, séances espacées de 3-7 jours
        3. Prothèses → Matinée, séances espacées de 14-21 jours
        4. Soins de routine → Après-midi, planning flexible
        5. Patients âgés → Matinée
        6. Enfants → Après-midi
        7. Adultes actifs → Début de matinée ou fin de journée

        RÉPONSE REQUISE:
        Fournis tes recommandations au format JSON avec:
        - timing_recommendations: pour chaque traitement
        - spacing_adjustments: modifications des délais
        - priority_notes: notes importantes
        - schedule_rationale: justification du planning
        """
        
        user_message = f"""
        CONTEXTE DE PLANIFICATION:
        
        Patient: {context['patient_info']['name']}
        Préférences patient: {context['patient_info']['preferences']}
        
        Traitements à planifier:
        {json.dumps(context['treatment_analysis'], indent=2, ensure_ascii=False)}
        
        Date de début souhaitée: {context['start_date']}
        
        Préférences du dentiste:
        - Jours de travail: {context['dentist_preferences']['working_days']}
        - Horaires: {context['dentist_preferences']['working_hours']}
        - Chirurgies: {context['dentist_preferences']['preferred_schedules']['surgeries']}
        - Consultations: {context['dentist_preferences']['preferred_schedules']['consultations']}
        
        Analyse ces traitements et propose un planning optimal avec justification.
        """
        
        try:
            response = self.client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_message}
                ],
                temperature=0.3,
                max_tokens=1500
            )
            
            # Try to parse JSON response
            content = response.choices[0].message.content
            
            # Extract JSON from response if it's wrapped in markdown
            if "```json" in content:
                json_start = content.find("```json") + 7
                json_end = content.find("```", json_start)
                json_content = content[json_start:json_end].strip()
            else:
                json_content = content
            
            try:
                llm_recommendations = json.loads(json_content)
            except json.JSONDecodeError:
                # Fallback to basic parsing if JSON parsing fails
                llm_recommendations = {
                    "timing_recommendations": [],
                    "spacing_adjustments": [],
                    "priority_notes": ["LLM response parsing failed - using default recommendations"],
                    "schedule_rationale": content
                }
            
            return llm_recommendations
            
        except Exception as e:
            print(f"❌ LLM scheduling error: {e}")
            return {
                "timing_recommendations": [],
                "spacing_adjustments": [],
                "priority_notes": [f"LLM error: {str(e)}"],
                "schedule_rationale": "Using fallback scheduling logic"
            }
    
    def apply_llm_recommendations(self, treatment_sequence: list, llm_response: dict, 
                                start_date: str, patient_id: str) -> dict:
        """Apply LLM recommendations to create optimized schedule"""
        
        try:
            print(f"🔍 DEBUG: Starting apply_llm_recommendations")
            print(f"🔍 DEBUG: treatment_sequence = {len(treatment_sequence)} treatments")
            print(f"🔍 DEBUG: llm_response keys = {list(llm_response.keys())}")
            print(f"🔍 DEBUG: start_date = {start_date}")
            print(f"🔍 DEBUG: patient_id = {patient_id}")
            
            optimized_appointments = []
            current_date = datetime.strptime(start_date, '%Y-%m-%d')
            
            for i, treatment in enumerate(treatment_sequence):
                print(f"🔍 DEBUG: Processing treatment {i}: {treatment}")
                
                treatment_class = self.classify_treatment(treatment.get('traitement', ''))
                print(f"🔍 DEBUG: treatment_class = {treatment_class} (type: {type(treatment_class)})")
                
                # Apply LLM timing recommendations if available
                preferred_time = self.get_optimal_time_for_treatment(
                    treatment, treatment_class, llm_response, i
                )
                print(f"🔍 DEBUG: preferred_time = {preferred_time}")
                
                # Apply intelligent date scheduling
                appointment_date = self.get_optimal_date_for_treatment(
                    current_date, treatment, treatment_class, llm_response, i
                )
                print(f"🔍 DEBUG: appointment_date = {appointment_date}")
                
                # Get available slots for the optimal date
                duration_minutes = self.parse_duration_minutes(treatment.get('duree', '60 min'))
                print(f"🔍 DEBUG: duration_minutes = {duration_minutes}")
                
                available_slots = self.practice_db.get_available_slots(
                    appointment_date.strftime('%Y-%m-%d'), 
                    duration_minutes
                )
                print(f"🔍 DEBUG: available_slots = {available_slots}")
                
                # Find best time slot
                final_time = self.find_best_time_slot(
                    available_slots, preferred_time, treatment_class
                )
                print(f"🔍 DEBUG: final_time = {final_time}")
                
                # If no slots available, try next working day
                if not final_time:
                    print(f"🔍 DEBUG: No slots available, finding next available slot")
                    appointment_date, final_time = self.find_next_available_slot(
                        appointment_date, duration_minutes, treatment_class
                    )
                    print(f"🔍 DEBUG: Next available: {appointment_date}, {final_time}")
                
                # Get scheduling reasoning
                reasoning = self.get_scheduling_reasoning(treatment, treatment_class, llm_response, i)
                print(f"🔍 DEBUG: reasoning = {reasoning}")
                
                optimized_appointments.append({
                    'date': appointment_date,
                    'time': final_time,
                    'treatment': treatment,
                    'duration_minutes': duration_minutes,
                    'classification': treatment_class,
                    'reasoning': reasoning
                })
                
                # Update current_date for next appointment based on treatment spacing
                current_date = self.calculate_next_appointment_date(
                    appointment_date, treatment, treatment_class, llm_response
                )
                print(f"🔍 DEBUG: Next current_date = {current_date}")
            
            print(f"🔍 DEBUG: Generating scheduling summary")
            scheduling_summary = self.generate_scheduling_summary(optimized_appointments, llm_response)
            
            result = {
                'appointments': optimized_appointments,
                'llm_analysis': llm_response,
                'total_duration': sum(apt['duration_minutes'] for apt in optimized_appointments),
                'scheduling_summary': scheduling_summary
            }
            
            print(f"🔍 DEBUG: apply_llm_recommendations completed successfully")
            return result
            
        except Exception as e:
            print(f"❌ ERROR in apply_llm_recommendations: {str(e)}")
            print(f"❌ ERROR type: {type(e)}")
            import traceback
            traceback.print_exc()
            raise e
    
    def get_optimal_time_for_treatment(self, treatment: dict, classification: dict, 
                                     llm_response: dict, step_index: int) -> str:
        """Determine optimal time for a specific treatment"""
        
        # Type check and fix classification if it's not a dict
        if not isinstance(classification, dict):
            print(f"⚠️  Warning: classification is not a dict, got {type(classification)}: {classification}")
            return '09:00'  # Default time
        
        # Check LLM recommendations first
        if llm_response.get('timing_recommendations'):
            for rec in llm_response['timing_recommendations']:
                # Type check - skip if rec is not a dict
                if not isinstance(rec, dict):
                    print(f"⚠️  Warning: timing_recommendation is not a dict, got {type(rec)}: {rec}")
                    continue
                if rec.get('step') == step_index + 1:
                    return rec.get('recommended_time', '09:00')
        
        # Use classification-based logic
        preferred_time = classification.get('preferred_time', 'morning')
        
        if preferred_time == 'morning':
            return '09:00'
        elif preferred_time == 'afternoon':
            return '14:00'
        else:
            return '09:00'  # Default
    
    def get_optimal_date_for_treatment(self, current_date: datetime, treatment: dict, 
                                     classification: dict, llm_response: dict, step_index: int) -> datetime:
        """Determine optimal date for a specific treatment"""
        
        # Type check and fix classification if it's not a dict
        if not isinstance(classification, dict):
            print(f"⚠️  Warning: classification is not a dict, got {type(classification)}: {classification}")
            # Return default 7 days later with working day adjustment
            target_date = current_date + timedelta(days=7)
            return self.adjust_to_working_day(target_date, {
                "category": "routine_treatments",
                "preferred_time": "afternoon",
                "duration_buffer": 10,
                "special_requirements": {}
            })
        
        # Apply LLM spacing adjustments if available
        if llm_response.get('spacing_adjustments'):
            for adj in llm_response['spacing_adjustments']:
                # Type check - skip if adj is not a dict
                if not isinstance(adj, dict):
                    print(f"⚠️  Warning: spacing_adjustment is not a dict, got {type(adj)}: {adj}")
                    continue
                if adj.get('step') == step_index + 1:
                    recommended_days = adj.get('recommended_days', 7)
                    target_date = current_date + timedelta(days=recommended_days)
                    return self.adjust_to_working_day(target_date, classification)
        
        # Use default spacing logic
        delay_str = treatment.get('delai', '1 semaine')
        delay_days = self.parse_delay_to_days(delay_str)
        
        # Apply treatment-specific spacing rules
        if classification.get('category') == 'surgical_treatments':
            delay_days = max(delay_days, 7)  # Minimum 1 week for surgery
        elif classification.get('category') == 'endodontic_treatments':
            delay_days = max(delay_days, 3)  # Minimum 3 days for endodontic
        
        target_date = current_date + timedelta(days=delay_days)
        return self.adjust_to_working_day(target_date, classification)
    
    def adjust_to_working_day(self, date: datetime, classification: dict) -> datetime:
        """Adjust date to working day considering treatment-specific rules"""
        
        # Skip weekends
        while date.weekday() >= 5:  # Saturday = 5, Sunday = 6
            date += timedelta(days=1)
        
        # Type check and fix classification if it's not a dict
        if not isinstance(classification, dict):
            print(f"⚠️  Warning: classification is not a dict, got {type(classification)}: {classification}")
            # Create a default classification dict
            classification = {
                "category": "routine_treatments",
                "preferred_time": "afternoon",
                "duration_buffer": 10,
                "special_requirements": {}
            }
        
        # Special rule: avoid Friday for surgical treatments
        try:
            if (classification.get('category') == 'surgical_treatments' and 
                classification.get('special_requirements', {}).get('avoid_friday', False)):
                if date.weekday() == 4:  # Friday
                    date += timedelta(days=3)  # Move to Monday
        except Exception as e:
            print(f"⚠️  Warning in adjust_to_working_day: {e}")
        
        return date
    
    def find_best_time_slot(self, available_slots: list, preferred_time: str, 
                           classification: dict) -> str:
        """Find the best available time slot"""
        
        if not available_slots:
            return None
        
        # Type check and fix classification if it's not a dict
        if not isinstance(classification, dict):
            print(f"⚠️  Warning: classification is not a dict, got {type(classification)}: {classification}")
            # Just return the first available slot or preferred time if available
            if preferred_time in available_slots:
                return preferred_time
            else:
                return available_slots[0]
        
        # If preferred time is available, use it
        if preferred_time in available_slots:
            return preferred_time
        
        # Find closest available slot to preferred time
        preferred_minutes = self.time_to_minutes(preferred_time)
        
        best_slot = min(available_slots, key=lambda slot: 
            abs(self.time_to_minutes(slot) - preferred_minutes))
        
        return best_slot
    
    def find_next_available_slot(self, start_date: datetime, duration_minutes: int, 
                               classification: dict) -> tuple:
        """Find next available slot if current date is full"""
        
        # Type check and fix classification if it's not a dict
        if not isinstance(classification, dict):
            print(f"⚠️  Warning: classification is not a dict, got {type(classification)}: {classification}")
            # Create a default classification dict
            classification = {
                "category": "routine_treatments",
                "preferred_time": "afternoon",
                "duration_buffer": 10,
                "special_requirements": {}
            }
        
        max_attempts = 14  # Try up to 2 weeks
        current_date = start_date
        
        for _ in range(max_attempts):
            current_date += timedelta(days=1)
            current_date = self.adjust_to_working_day(current_date, classification)
            
            available_slots = self.practice_db.get_available_slots(
                current_date.strftime('%Y-%m-%d'), 
                duration_minutes
            )
            
            if available_slots:
                # Get preferred time for this treatment type
                preferred_time = self.get_preferred_time_for_classification(classification)
                best_slot = self.find_best_time_slot(available_slots, preferred_time, classification)
                return current_date, best_slot
        
        # If no slots found, return original date with default time
        return start_date, '09:00'
    
    def get_preferred_time_for_classification(self, classification: dict) -> str:
        """Get preferred time based on treatment classification"""
        # Type check and fix classification if it's not a dict
        if not isinstance(classification, dict):
            print(f"⚠️  Warning: classification is not a dict, got {type(classification)}: {classification}")
            return '09:00'  # Default time
        
        preferred = classification.get('preferred_time', 'morning')
        
        if preferred == 'morning':
            return '09:00'
        elif preferred == 'afternoon':
            return '14:00'
        else:
            return '09:00'
    
    def parse_duration_minutes(self, duration_str: str) -> int:
        """Parse duration string to minutes"""
        try:
            if 'min' in duration_str:
                return int(duration_str.split('min')[0].strip())
            elif 'h' in duration_str:
                hours = float(duration_str.split('h')[0].strip())
                return int(hours * 60)
            else:
                return 60  # Default
        except:
            return 60
    
    def parse_delay_to_days(self, delay_str: str) -> int:
        """Parse delay string to days"""
        try:
            if 'jour' in delay_str:
                return int(delay_str.split('jour')[0].strip())
            elif 'semaine' in delay_str:
                weeks = int(delay_str.split('semaine')[0].strip())
                return weeks * 7
            elif 'mois' in delay_str:
                months = int(delay_str.split('mois')[0].strip())
                return months * 30
            else:
                return 7  # Default 1 week
        except:
            return 7
    
    def time_to_minutes(self, time_str: str) -> int:
        """Convert time string to minutes since midnight"""
        try:
            hours, minutes = map(int, time_str.split(':'))
            return hours * 60 + minutes
        except:
            return 540  # Default to 9:00 AM
    
    def calculate_next_appointment_date(self, current_date: datetime, treatment: dict, 
                                      classification: dict, llm_response: dict) -> datetime:
        """Calculate the base date for the next appointment"""
        
        # Type check and fix classification if it's not a dict
        if not isinstance(classification, dict):
            print(f"⚠️  Warning: classification is not a dict, got {type(classification)}: {classification}")
            # Return default 7 days later
            return current_date + timedelta(days=7)
        
        # Get treatment-specific spacing
        spacing_rules = self.dentist_preferences.get('treatment_spacing', {})
        
        if classification.get('category') == 'surgical_treatments':
            min_days = spacing_rules.get('surgical', {}).get('min_days', 7)
        elif classification.get('category') == 'endodontic_treatments':
            min_days = spacing_rules.get('endodontic', {}).get('min_days', 3)
        elif classification.get('category') == 'prosthetic_treatments':
            min_days = spacing_rules.get('prosthetic', {}).get('min_days', 14)
        else:
            min_days = spacing_rules.get('routine', {}).get('min_days', 7)
        
        return current_date + timedelta(days=min_days)
    
    def get_scheduling_reasoning(self, treatment: dict, classification: dict, 
                               llm_response: dict, step_index: int) -> str:
        """Generate reasoning for scheduling decision"""
        
        # Type check and fix classification if it's not a dict
        if not isinstance(classification, dict):
            print(f"⚠️  Warning: classification is not a dict, got {type(classification)}: {classification}")
            return "Programmation standard (classification invalide)"
        
        reasons = []
        
        # Treatment-based reasoning
        if classification.get('category') == 'surgical_treatments':
            reasons.append("Chirurgie programmée le matin pour optimiser la récupération")
        elif classification.get('category') == 'endodontic_treatments':
            reasons.append("Traitement endodontique programmé le matin pour concentration maximale")
        elif classification.get('category') == 'prosthetic_treatments':
            reasons.append("Travail prothétique programmé le matin pour précision optimale")
        
        # LLM reasoning if available
        if llm_response.get('priority_notes'):
            for note in llm_response['priority_notes']:
                # Type check - ensure note is a string
                if isinstance(note, str) and str(step_index + 1) in note:
                    reasons.append(f"IA: {note}")
                elif not isinstance(note, str):
                    print(f"⚠️  Warning: priority_note is not a string, got {type(note)}: {note}")
        
        return " | ".join(reasons) if reasons else "Programmation standard"
    
    def generate_scheduling_summary(self, appointments: list, llm_response: dict) -> dict:
        """Generate a summary of the scheduling decisions"""
        
        summary = {
            'total_appointments': len(appointments),
            'morning_appointments': sum(1 for apt in appointments if apt['time'] < '12:00'),
            'afternoon_appointments': sum(1 for apt in appointments if apt['time'] >= '12:00'),
            'surgical_appointments': 0,
            'routine_appointments': 0,
            'ai_recommendations_applied': len(llm_response.get('timing_recommendations', [])),
            'scheduling_rationale': llm_response.get('schedule_rationale', 'Planification intelligente appliquée'),
            'insights': []
        }
        
        # Count surgical and routine appointments with error handling
        for apt in appointments:
            classification = apt.get('classification', {})
            if isinstance(classification, dict):
                if classification.get('category') == 'surgical_treatments':
                    summary['surgical_appointments'] += 1
                elif classification.get('category') == 'routine_treatments':
                    summary['routine_appointments'] += 1
            else:
                print(f"⚠️  Warning: classification is not a dict in summary, got {type(classification)}: {classification}")
                summary['routine_appointments'] += 1  # Default to routine
        
        # Add LLM insights if available
        if llm_response.get('summary_points'):
            for point in llm_response['summary_points']:
                # Type check - ensure point is a string
                if isinstance(point, str):
                    summary['insights'].append(point)
                else:
                    print(f"⚠️  Warning: summary_point is not a string, got {type(point)}: {point}")
        
        return summary

# Initialize the intelligent scheduler
intelligent_scheduler = IntelligentScheduler(client, practice_db)

def get_current_schedule_context():
    """Get current schedule context for the LLM"""
    try:
        # Get current week appointments
        today = datetime.now()
        week_start = today - timedelta(days=today.weekday())
        week_end = week_start + timedelta(days=6)
        
        appointments = practice_db.get_appointments_by_date_range(
            week_start.strftime('%Y-%m-%d'),
            week_end.strftime('%Y-%m-%d')
        )
        
        # Get practice preferences
        dentist_prefs = intelligent_scheduler.get_dentist_preferences()
        
        # Build context
        context = f"""
SEMAINE ACTUELLE: {week_start.strftime('%d/%m/%Y')} - {week_end.strftime('%d/%m/%Y')}

RENDEZ-VOUS PROGRAMMÉS:
"""
        
        # Group appointments by day
        days_appointments = {}
        for apt in appointments:
            apt_date = apt['date']
            if apt_date not in days_appointments:
                days_appointments[apt_date] = []
            days_appointments[apt_date].append(apt)
        
        # Add appointments to context
        for date in sorted(days_appointments.keys()):
            day_name = datetime.strptime(date, '%Y-%m-%d').strftime('%A %d/%m')
            context += f"\n{day_name}:\n"
            
            for apt in sorted(days_appointments[date], key=lambda x: x['time']):
                patient = practice_db.get_patient(apt['patient_id'])
                patient_name = f"{patient.get('first_name', '')} {patient.get('last_name', '')}" if patient else "Patient inconnu"
                context += f"  - {apt['time']}: {patient_name} - {apt['treatment']} ({apt['duration']})\n"
        
        # Add available slots information
        context += f"\n\nPRÉFÉRENCES PRATICIEN:\n"
        context += f"- Jours de travail: {', '.join(dentist_prefs['working_days'])}\n"
        context += f"- Horaires: {dentist_prefs['working_hours']['start']} - {dentist_prefs['working_hours']['end']}\n"
        context += f"- Pause déjeuner: {dentist_prefs['working_hours']['lunch_break']['start']} - {dentist_prefs['working_hours']['lunch_break']['end']}\n"
        
        return context
        
    except Exception as e:
        print(f"Error getting schedule context: {e}")
        return "Erreur lors du chargement du contexte planning"

def parse_schedule_actions(llm_response):
    """Parse LLM response for actionable schedule items"""
    try:
        actions = []
        
        # Look for common action patterns
        response_lower = llm_response.lower()
        
        # Reschedule patterns
        if any(word in response_lower for word in ['reprogrammer', 'déplacer', 'changer', 'modifier']):
            actions.append({
                'type': 'reschedule',
                'description': 'Reprogrammation de rendez-vous détectée',
                'priority': 'high'
            })
        
        # Find slot patterns
        if any(word in response_lower for word in ['trouver', 'créer', 'libérer', 'créneau']):
            actions.append({
                'type': 'find_slot',
                'description': 'Recherche de créneau détectée',
                'priority': 'medium'
            })
        
        # Optimize patterns
        if any(word in response_lower for word in ['optimiser', 'améliorer', 'organiser']):
            actions.append({
                'type': 'optimize',
                'description': 'Optimisation du planning détectée',
                'priority': 'low'
            })
        
        # Emergency patterns
        if any(word in response_lower for word in ['urgence', 'urgent', 'immédiat']):
            actions.append({
                'type': 'emergency',
                'description': 'Gestion d\'urgence détectée',
                'priority': 'urgent'
            })
        
        return actions
        
    except Exception as e:
        print(f"Error parsing schedule actions: {e}")
        return []

def analyze_schedule_request(user_message, schedule_context):
    """Analyze user request and extract specific schedule actions"""
    try:
        print(f"🔍 Analyzing schedule request: {user_message}")
        
        # Enhanced analysis prompt with French date detection
        analysis_prompt = f"""
Tu es un assistant intelligent de planification dentaire. Analyse cette demande du DENTISTE et propose des actions concrètes.

CONTEXTE PLANNING:
{schedule_context}

DEMANDE DU DENTISTE:
{user_message}

INSTRUCTIONS IMPORTANTES:
1. Le dentiste demande des modifications à SON planning
2. Détecte les dates françaises (format DD/MM, "vendredi 10/7", "lundi prochain", etc.)
3. Propose des actions concrètes avec confirmation DENTISTE (pas patient)
4. Fournis des solutions alternatives si besoin
5. Pour les reprogrammations, target_date doit être la date AVEC les RDV à reprogrammer

FORMATS DE DATES À DÉTECTER:
- "vendredi 10/7" → 10/07/2024
- "lundi 15" → 15 du mois courant
- "demain" → date de demain
- "la semaine prochaine" → dates de la semaine suivante

Réponds au format JSON avec:
{{
    "analysis": "Analyse de la demande du dentiste",
    "detected_dates": ["dates détectées au format YYYY-MM-DD"],
    "proposed_actions": [
        {{
            "type": "reschedule|find_slot|block_time|emergency",
            "description": "Description de l'action",
            "priority": "urgent|high|medium|low",
            "target_date": "date des RDV à reprogrammer (même date que detected_dates)",
            "reason": "raison de la modification (absence, urgence, etc.)",
            "requires_dentist_approval": true,
            "proposed_solutions": ["Solutions concrètes proposées"]
        }}
    ],
    "immediate_actions": ["Actions à effectuer immédiatement"],
    "confirmation_needed": "Message de confirmation pour le dentiste"
}}
"""
        
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "Tu es un assistant de planification dentaire qui analyse les demandes du dentiste et propose des actions concrètes avec confirmation dentiste."},
                {"role": "user", "content": analysis_prompt}
            ],
            temperature=0.3,
            max_tokens=1200
        )
        
        content = response.choices[0].message.content
        print(f"🔍 LLM Analysis response: {content}")
        
        # Extract JSON from response
        if "```json" in content:
            json_start = content.find("```json") + 7
            json_end = content.find("```", json_start)
            json_content = content[json_start:json_end].strip()
        else:
            json_content = content
        
        try:
            analysis = json.loads(json_content)
            print(f"🔍 Parsed analysis: {analysis}")
            
            # Enhance with date detection if not already present
            if not analysis.get('detected_dates'):
                analysis['detected_dates'] = extract_dates_from_french_text(user_message)
            
            return analysis
            
        except json.JSONDecodeError as e:
            print(f"❌ JSON parsing error: {e}")
            print(f"Raw content: {content}")
            
            # Fallback analysis
            return {
                "analysis": "Analyse automatique de la demande",
                "detected_dates": extract_dates_from_french_text(user_message),
                "proposed_actions": [{
                    "type": "reschedule",
                    "description": "Reprogrammation nécessaire",
                    "priority": "high",
                    "target_date": None,
                    "reason": "Demande de modification du planning",
                    "requires_dentist_approval": True,
                    "proposed_solutions": ["Recherche de créneaux alternatifs"]
                }],
                "immediate_actions": ["Identifier les rendez-vous concernés"],
                "confirmation_needed": "Confirmez-vous cette modification de planning ?",
                "raw_response": content
            }
            
    except Exception as e:
        print(f"❌ Error analyzing schedule request: {e}")
        return {
            "analysis": f"Erreur d'analyse: {str(e)}",
            "detected_dates": [],
            "proposed_actions": [],
            "immediate_actions": [],
            "confirmation_needed": "Une erreur s'est produite lors de l'analyse."
        }

def extract_dates_from_french_text(text):
    """Extract dates from French text"""
    import re
    
    dates = []
    text_lower = text.lower()
    
    # Pattern for DD/MM format
    date_pattern = r'\b(\d{1,2})[\/\-](\d{1,2})(?:[\/\-](\d{4}))?\b'
    matches = re.findall(date_pattern, text)
    
    for match in matches:
        day, month, year = match
        if not year:
            year = datetime.now().year
        
        try:
            date_obj = datetime(int(year), int(month), int(day))
            dates.append(date_obj.strftime('%Y-%m-%d'))
        except ValueError:
            continue
    
    # Handle relative dates
    today = datetime.now()
    
    if 'demain' in text_lower:
        dates.append((today + timedelta(days=1)).strftime('%Y-%m-%d'))
    
    if 'après-demain' in text_lower:
        dates.append((today + timedelta(days=2)).strftime('%Y-%m-%d'))
    
    # Handle day names (vendredi, lundi, etc.)
    days_fr = {
        'lundi': 0, 'mardi': 1, 'mercredi': 2, 'jeudi': 3, 
        'vendredi': 4, 'samedi': 5, 'dimanche': 6
    }
    
    for day_name, day_num in days_fr.items():
        if day_name in text_lower:
            # Find next occurrence of this day
            days_ahead = (day_num - today.weekday()) % 7
            if days_ahead == 0:  # Today
                days_ahead = 7  # Next week
            target_date = today + timedelta(days=days_ahead)
            dates.append(target_date.strftime('%Y-%m-%d'))
    
    return list(set(dates))  # Remove duplicates

def find_appointments_for_date(date_str):
    """Find appointments for a specific date"""
    try:
        print(f"🔍 Finding appointments for date: {date_str}")
        
        # Parse date string (handle various formats)
        if '/' in date_str:
            # Handle DD/MM or DD/MM/YYYY format
            parts = date_str.split('/')
            if len(parts) == 2:
                day, month = parts
                year = datetime.now().year
            elif len(parts) == 3:
                day, month, year = parts
            else:
                raise ValueError(f"Invalid date format: {date_str}")
            
            # Convert to YYYY-MM-DD format
            formatted_date = f"{year}-{month.zfill(2)}-{day.zfill(2)}"
        elif '-' in date_str:
            # Already in YYYY-MM-DD format
            formatted_date = date_str
        else:
            # Try to parse as DD/MM without separator
            if len(date_str) == 4:  # DDMM
                day = date_str[:2]
                month = date_str[2:]
                year = datetime.now().year
                formatted_date = f"{year}-{month.zfill(2)}-{day.zfill(2)}"
            else:
                raise ValueError(f"Cannot parse date: {date_str}")
        
        print(f"🔍 Formatted date: {formatted_date}")
        
        # Get appointments from database
        appointments = practice_db.get_appointments_for_date(formatted_date)
        print(f"🔍 Found {len(appointments)} appointments")
        
        # Format appointments for frontend
        formatted_appointments = []
        for apt in appointments:
            formatted_appointments.append({
                'id': apt['id'],
                'patient_id': apt['patient_id'],
                'date': apt['appointment_date'],
                'time': apt['appointment_time'],
                'duration_minutes': apt.get('duration_minutes', 60),
                'treatment': apt.get('treatment_type', 'Consultation'),
                'patient_name': f"{apt.get('first_name', '')} {apt.get('last_name', '')}".strip(),
                'doctor': apt.get('doctor', 'Dr.'),
                'status': apt.get('status', 'scheduled')
            })
        
        return formatted_appointments
        
    except Exception as e:
        print(f"❌ Error finding appointments for date {date_str}: {e}")
        return []

def propose_reschedule_options(appointments, target_date=None):
    """Propose intelligent reschedule options using AI analysis"""
    try:
        print(f"🧠 AI analyzing reschedule options for {len(appointments)} appointments")
        
        if not appointments:
            return []
        
        # Get available slots for the next 2 weeks
        today = datetime.now()
        available_slots_by_date = {}
        
        for i in range(14):  # Next 2 weeks
            date = today + timedelta(days=i)
            if date.weekday() < 5:  # Monday to Friday only
                date_str = date.strftime('%Y-%m-%d')
                day_name = date.strftime('%A')
                french_day_name = {
                    'Monday': 'Lundi',
                    'Tuesday': 'Mardi', 
                    'Wednesday': 'Mercredi',
                    'Thursday': 'Jeudi',
                    'Friday': 'Vendredi'
                }.get(day_name, day_name)
                
                slots = practice_db.get_available_slots(date_str, 60)
                if slots:
                    available_slots_by_date[date_str] = {
                        'date': date_str,
                        'day_name': f"{french_day_name} {date.strftime('%d/%m')}",
                        'slots': slots[:6]  # Limit to 6 slots per day
                    }
        
        # Format reschedule options for frontend
        reschedule_options = []
        
        for apt in appointments:
            patient_name = apt.get('patient_name', 'Patient inconnu')
            if not patient_name or patient_name.strip() == '':
                # Try to construct from patient data
                patient_id = apt.get('patient_id')
                if patient_id:
                    patient = practice_db.get_patient(patient_id)
                    if patient:
                        patient_name = f"{patient.get('first_name', '')} {patient.get('last_name', '')}".strip()
                
                if not patient_name:
                    patient_name = 'Patient inconnu'
            
            # Convert available slots to the expected format
            available_options = []
            for date_str, date_info in available_slots_by_date.items():
                if date_info['slots']:  # Only include dates with available slots
                    available_options.append({
                        'date': date_str,
                        'day_name': date_info['day_name'],
                        'slots': date_info['slots']
                    })
            
            reschedule_option = {
                'appointment_id': apt.get('id'),
                'patient_id': apt.get('patient_id'),
                'patient_name': patient_name,
                'treatment': apt.get('treatment', 'Traitement'),
                'current_date': apt.get('date'),
                'current_time': apt.get('time'),
                'duration_minutes': apt.get('duration_minutes', 60),
                'available_options': available_options
            }
            
            reschedule_options.append(reschedule_option)
        
        print(f"🔍 Generated {len(reschedule_options)} reschedule options")
        return reschedule_options
        
    except Exception as e:
        print(f"❌ Error in propose_reschedule_options: {e}")
        import traceback
        traceback.print_exc()
        return []

        return jsonify({'error': f'Erreur du serveur: {str(e)}'}), 500

@app.route('/api/schedule-execute-action', methods=['POST'])
def execute_schedule_action():
    """Execute intelligent schedule actions automatically"""
    try:
        data = request.json
        action = data.get('action')
        
        print(f"🚀 Executing intelligent schedule action: {action}")
        
        if action == 'reschedule':
            appointments = data.get('appointments', [])
            
            if not appointments:
                return jsonify({
                    'success': False,
                    'error': 'Aucun rendez-vous trouvé à reprogrammer'
                })
            
            # Use intelligent rescheduling
            execution_results = propose_reschedule_options(appointments)
            
            # Format response for frontend
            if execution_results:
                successful_reschedules = [r for r in execution_results if r.get('success')]
                failed_reschedules = [r for r in execution_results if not r.get('success')]
                
                response_message = "🤖 **Reprogrammation intelligente terminée !**\n\n"
                
                if successful_reschedules:
                    response_message += "✅ **Rendez-vous reprogrammés automatiquement :**\n"
                    for result in successful_reschedules:
                        response_message += f"• **{result['patient_name']}** : {result['old_slot']} → **{result['new_slot']}**\n"
                        response_message += f"  💭 *{result['reasoning']}*\n"
                        response_message += f"  🎯 *Confiance : {result['confidence']:.1%}*\n\n"
                
                if failed_reschedules:
                    response_message += "⚠️ **Rendez-vous nécessitant attention manuelle :**\n"
                    for result in failed_reschedules:
                        response_message += f"• **{result['decision']['patient_name']}** : {result.get('error', 'Erreur inconnue')}\n"
                
                response_message += f"\n📊 **Bilan :** {len(successful_reschedules)} reprogrammés automatiquement, {len(failed_reschedules)} nécessitent votre attention."
                
                return jsonify({
                    'success': True,
                    'message': response_message,
                    'execution_results': execution_results,
                    'stats': {
                        'total': len(execution_results),
                        'successful': len(successful_reschedules),
                        'failed': len(failed_reschedules)
                    }
                })
            else:
                return jsonify({
                    'success': False,
                    'message': "❌ Aucune reprogrammation automatique n'a pu être effectuée. Veuillez vérifier les créneaux disponibles."
                })
        
        elif action == 'block_time':
            date = data.get('date')
            time_range = data.get('time_range', {})
            reason = data.get('reason', 'Temps bloqué')
            
            # Block time slots
            success = block_time_slots(date, time_range, reason)
            
            if success:
                return jsonify({
                    'success': True,
                    'message': f"🚫 Créneaux bloqués pour le {date} : {time_range.get('start', '')} - {time_range.get('end', '')}"
                })
            else:
                return jsonify({
                    'success': False,
                    'message': "❌ Erreur lors du blocage des créneaux"
                })
        
        elif action == 'emergency_slot':
            date = data.get('date')
            time = data.get('time')
            patient_info = data.get('patient_info', {})
            
            # Create emergency slot
            success = create_emergency_slot(date, time, patient_info)
            
            if success:
                return jsonify({
                    'success': True,
                    'message': f"🚨 Créneau d'urgence créé pour le {date} à {time}"
                })
            else:
                return jsonify({
                    'success': False,
                    'message': "❌ Erreur lors de la création du créneau d'urgence"
                })
        
        else:
            return jsonify({
                'success': False,
                'error': f'Action non reconnue: {action}'
            })
            
    except Exception as e:
        print(f"❌ Error executing schedule action: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({
            'success': False,
            'error': str(e)
        })

def block_time_slots(date, time_range, reason):
    """Block time slots for a specific date and time range"""
    try:
        start_time = time_range.get('start', '09:00')
        end_time = time_range.get('end', '17:00')
        
        # Create blocking appointment
        blocking_appointment = {
            'date': date,
            'time': start_time,
            'patient_name': 'BLOQUÉ',
            'treatment': reason,
            'duration': calculate_duration(start_time, end_time),
            'status': 'blocked',
            'type': 'block'
        }
        
        success = practice_db.add_appointment(blocking_appointment)
        print(f"🚫 Blocked time slots: {date} {start_time}-{end_time}")
        
        return success
        
    except Exception as e:
        print(f"❌ Error blocking time slots: {e}")
        return False

def create_emergency_slot(date, time, patient_info):
    """Create an emergency appointment slot"""
    try:
        emergency_appointment = {
            'date': date,
            'time': time,
            'patient_name': patient_info.get('name', 'Urgence'),
            'treatment': 'Urgence dentaire',
            'duration': 30,
            'status': 'emergency',
            'type': 'emergency',
            'phone': patient_info.get('phone', ''),
            'notes': patient_info.get('notes', 'Rendez-vous d\'urgence')
        }
        
        success = practice_db.add_appointment(emergency_appointment)
        print(f"🚨 Created emergency slot: {date} {time}")
        
        return success
        
    except Exception as e:
        print(f"❌ Error creating emergency slot: {e}")
        return False

def calculate_duration(start_time, end_time):
    """Calculate duration between two times in minutes"""
    try:
        from datetime import datetime
        
        start = datetime.strptime(start_time, '%H:%M')
        end = datetime.strptime(end_time, '%H:%M')
        
        duration = (end - start).total_seconds() / 60
        return int(duration)
        
    except Exception as e:
        print(f"❌ Error calculating duration: {e}")
        return 60  # Default to 1 hour

@app.route('/api/execute-autonomous-plan', methods=['POST'])
def execute_autonomous_plan():
    """Execute an approved autonomous rescheduling plan"""
    try:
        data = request.get_json()
        autonomous_plan = data.get('autonomous_plan')
        
        if not autonomous_plan:
            return jsonify({'error': 'Plan autonome manquant'}), 400
        
        if not autonomous_plan.get('execution_ready', False):
            return jsonify({'error': 'Plan non prêt pour exécution'}), 400
        
        decisions = autonomous_plan.get('decisions', [])
        if not decisions:
            return jsonify({'error': 'Aucune décision à exécuter'}), 400
        
        print(f"🚀 Executing autonomous plan with {len(decisions)} decisions...")
        
        # Execute each decision
        execution_results = []
        successful_executions = 0
        failed_executions = 0
        
        for decision in decisions:
            if not decision.get('success', False):
                execution_results.append({
                    'appointment_id': decision.get('appointment_id'),
                    'success': False,
                    'message': f"Décision non exécutable: {decision.get('reasoning', 'Raison inconnue')}"
                })
                failed_executions += 1
                continue
            
            # Execute the rescheduling
            try:
                result = execute_single_reschedule(decision)
                execution_results.append(result)
                
                if result.get('success', False):
                    successful_executions += 1
                else:
                    failed_executions += 1
                    
            except Exception as e:
                print(f"❌ Error executing decision for appointment {decision.get('appointment_id')}: {e}")
                execution_results.append({
                    'appointment_id': decision.get('appointment_id'),
                    'success': False,
                    'message': f"Erreur lors de l'exécution: {str(e)}"
                })
                failed_executions += 1
        
        # Generate summary
        summary = {
            'total_decisions': len(decisions),
            'successful_executions': successful_executions,
            'failed_executions': failed_executions,
            'success_rate': (successful_executions / len(decisions)) * 100 if decisions else 0
        }
        
        # Generate response message
        if successful_executions == len(decisions):
            message = f"🎉 Plan exécuté avec succès! Tous les {successful_executions} rendez-vous ont été reprogrammés."
        elif successful_executions > 0:
            message = f"✅ Plan partiellement exécuté: {successful_executions}/{len(decisions)} rendez-vous reprogrammés avec succès."
        else:
            message = f"❌ Échec de l'exécution: Aucun rendez-vous n'a pu être reprogrammé."
        
        response_data = {
            'success': successful_executions > 0,
            'message': message,
            'summary': summary,
            'execution_results': execution_results
        }
        
        print(f"✅ Autonomous plan execution completed: {successful_executions}/{len(decisions)} successful")
        return jsonify(response_data)
        
    except Exception as e:
        print(f"❌ Error executing autonomous plan: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({
            'success': False,
            'error': f'Erreur lors de l\'exécution du plan: {str(e)}',
            'summary': {'total_decisions': 0, 'successful_executions': 0, 'failed_executions': 0, 'success_rate': 0},
            'execution_results': []
        }), 500

def execute_single_reschedule(decision):
    """Execute a single rescheduling decision"""
    try:
        appointment_id = decision.get('appointment_id')
        new_date = decision.get('new_date')
        new_time = decision.get('new_time')
        
        if not all([appointment_id, new_date, new_time]):
            return {
                'appointment_id': appointment_id,
                'success': False,
                'message': 'Données de reprogrammation incomplètes'
            }
        
        # Get the current appointment using the correct method
        appointments = practice_db.get_appointments()
        current_appointment = None
        for apt in appointments:
            if apt.get('id') == appointment_id:
                current_appointment = apt
                break
        
        if not current_appointment:
            return {
                'appointment_id': appointment_id,
                'success': False,
                'message': 'Rendez-vous non trouvé'
            }
        
        # Check if the new slot is available using the correct method
        duration_minutes = current_appointment.get('duration_minutes', 60)
        available_slots = practice_db.get_available_slots(new_date, duration_minutes)
        
        if new_time not in available_slots:
            return {
                'appointment_id': appointment_id,
                'success': False,
                'message': 'Créneau non disponible'
            }
        
        # Update the appointment using direct database access
        conn = sqlite3.connect(practice_db.db_path)
        cursor = conn.cursor()
        
        cursor.execute('''
            UPDATE appointments 
            SET appointment_date = ?, appointment_time = ?, updated_at = ?
            WHERE id = ?
        ''', (new_date, new_time, datetime.now().isoformat(), appointment_id))
        
        success = cursor.rowcount > 0
        conn.commit()
        conn.close()
        
        if success:
            # Log the change
            print(f"✅ Appointment {appointment_id} rescheduled from {current_appointment.get('appointment_date')} {current_appointment.get('appointment_time')} to {new_date} {new_time}")
            
            return {
                'appointment_id': appointment_id,
                'success': True,
                'message': f"Rendez-vous reprogrammé avec succès vers {new_date} à {new_time}",
                'old_slot': f"{current_appointment.get('appointment_date')} à {current_appointment.get('appointment_time')}",
                'new_slot': f"{new_date} à {new_time}"
            }
        else:
            return {
                'appointment_id': appointment_id,
                'success': False,
                'message': 'Échec de la mise à jour en base de données'
            }
            
    except Exception as e:
        print(f"❌ Error executing single reschedule: {e}")
        import traceback
        traceback.print_exc()
        return {
            'appointment_id': decision.get('appointment_id'),
            'success': False,
            'message': f'Erreur lors de la reprogrammation: {str(e)}'
        }

if __name__ == '__main__':
    print("🚀 Starting Enhanced Dental Practice Management System...")
    print("📊 Database initialized")
    print("🧠 RAG system ready")
    print("🔧 Multi-LLM architecture active")
    print("🏥 Practice management features enabled")
    print("🎯 PowerPoint generation system integrated")
    print("🌐 Server starting on http://localhost:5001")
    app.run(debug=True, host='0.0.0.0', port=5001) 